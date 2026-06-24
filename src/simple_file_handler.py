"""
Simple File Handler - Teams Attachment Processing
Handles both direct uploads and SharePoint/OneDrive files via Graph API.
Extracted content is returned to the caller for processing and caching.
"""

import io
import logging
import base64
import csv
import re
from typing import Optional
from datetime import datetime
import json
import os
import tempfile
from config import Config

logger = logging.getLogger(__name__)


def _smart_truncate(text: str, file_type: str) -> str:
    """Truncate very large extracted content while logging useful diagnostics.

    Uses Config.MAX_EXTRACTED_CHARS as a hard cap. If the content is below
    the limit it is returned unchanged; otherwise we keep most of it and
    append a clear truncation notice so the LLM/user understands why
    the tail of the document is missing.
    """
    if text is None:
        return ""

    text = str(text)

    if len(text) <= Config.MAX_EXTRACTED_CHARS:
        size_mb = len(text.encode("utf-8", errors="ignore")) / (1024 * 1024)
        logger.info(
            f"\u2713 {file_type} content extracted FULLY: {len(text):,} chars ({size_mb:.2f} MB)"
        )
        return text

    kept_chars = int(Config.MAX_EXTRACTED_CHARS * 0.95)
    truncated = text[:kept_chars]

    original_size_mb = len(text.encode("utf-8", errors="ignore")) / (1024 * 1024)
    kept_size_mb = len(truncated.encode("utf-8", errors="ignore")) / (1024 * 1024)

    truncation_notice = (
        "\n\n⚠️ **CONTENT TRUNCATED (EXTREMELY LARGE FILE)**\n"
        f"Original: {len(text):,} chars ({original_size_mb:.1f} MB)\n"
        f"Extracted: {len(truncated):,} chars ({kept_size_mb:.1f} MB)\n"
        "Reason: File exceeds 20M character limit\n\n"
        "💡 This file is exceptionally large. For best results:\n"
        "• Split into smaller files\n"
        "• Ask specific questions about sections\n"
        "• Request analysis of particular ranges"
    )

    logger.warning(
        f"Truncated {file_type} content (EXTREMELY LARGE): {len(text):,} -> {len(truncated):,} chars"
    )
    return truncated + truncation_notice


def _llm_safe_cap(text: str, max_chars: int = None) -> str:
    """Secondary LLM-safe cap applied AFTER extraction.

    Cache may store the full _smart_truncate result (up to MAX_EXTRACTED_CHARS),
    but anything entering an LLM prompt MUST be capped much lower.
    Default cap: Config.MAX_DOC_SNIPPET_CHARS (typically 6000 chars).
    """
    if not text:
        return ""
    if max_chars is None:
        max_chars = getattr(Config, "MAX_DOC_SNIPPET_CHARS", 6000)
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + "\n[...CONTENT CAPPED FOR LLM SAFETY...]"


# Optional document processing imports
pypdf = None
Document = None
load_workbook = None
Image = None
Presentation = None
xlrd = None
textract = None
try:
    import pypdf
except ImportError:
    pass
try:
    from docx import Document
except ImportError:
    pass
try:
    from openpyxl import load_workbook
except ImportError:
    pass
try:
    from PIL import Image
except ImportError:
    pass
try:
    from pptx import Presentation
except ImportError:
    pass
try:
    import xlrd
except ImportError:
    pass
try:
    import textract
except ImportError:
    pass


def _format_cell_value(cell) -> str:
    """Format a cell value for AI readability - NO truncation, full content preserved."""
    if cell is None:
        return ""
    if isinstance(cell, float):
        # Check if it's actually an integer
        if cell == int(cell):
            return str(int(cell))
        # Format decimals nicely (4 decimal places max)
        return f"{cell:.4f}".rstrip('0').rstrip('.')
    if isinstance(cell, (int, bool)):
        return str(cell)
    # Convert to string and clean - FULL content preserved
    cell_str = str(cell).strip()
    # Replace problematic characters that could break table structure
    cell_str = cell_str.replace('\n', ' ').replace('\r', ' ').replace('\t', ' ')
    # Collapse multiple spaces
    while '  ' in cell_str:
        cell_str = cell_str.replace('  ', ' ')
    return cell_str


def _is_numeric(value: str) -> bool:
    """Check if a string value is numeric."""
    if not value:
        return False
    try:
        float(value.replace(",", "").replace("$", "").replace("%", ""))
        return True
    except ValueError:
        return False


def _textract_legacy_office(content: bytes, extension: str, display_name: str) -> Optional[str]:
    """Extract text from legacy Office formats (.ppt, .doc, .xls) using textract if available."""
    if textract is None:
        return None
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{extension}") as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        extracted = textract.process(tmp_path, extension=extension)
        if extracted:
            return extracted.decode("utf-8", errors="ignore").strip()
    except Exception as e:
        logger.warning(f"Legacy extraction failed for {display_name}: {e}")
    finally:
        if tmp_path:
            try:
                os.remove(tmp_path)
            except Exception:
                pass
    return None


def _is_sharepoint_url(url: str) -> bool:
    """Check if URL is from SharePoint or OneDrive (requires Graph API)."""
    if not url:
        return False
    url_lower = url.lower()
    return "sharepoint.com" in url_lower or "onedrive.com" in url_lower


def _analyze_csv_content(csv_text: str, display_name: str) -> str:
    """
    Analyze CSV content intelligently:
    - Parse and summarize data
    - Calculate totals for numeric columns
    - Identify key metrics and groupings
    - Return structured insights for the LLM
    """
    try:
        import pandas as pd
        import io
        
        # Parse CSV
        df = pd.read_csv(io.StringIO(csv_text))
        
        # Get basic info
        rows, cols = df.shape
        column_names = df.columns.tolist()
        
        from data_calculator import is_identifier_column, measure_columns, categorical_columns

        # Find numeric measure columns and calculate totals. Numeric identifiers
        # such as ProjectID are labels, not values to add up.
        numeric_cols = measure_columns(df)
        identifier_cols = [
            col for col in df.select_dtypes(include=['number']).columns.tolist()
            if is_identifier_column(col, df[col])
        ]
        categorical_cols = categorical_columns(df)
        totals = {}
        
        for col in numeric_cols:
            total = df[col].sum()
            totals[col] = total
        
        # Build analysis report
        report = f"**CSV Analysis**: {display_name}\n"
        report += f"**Rows:** {rows:,} | **Columns:** {cols}\n\n"
        
        # Summary section
        report += "**Column Info:**\n"
        for col in column_names:
            if col in numeric_cols:
                total = totals[col]
                avg = df[col].mean()
                report += f"  • {col}: NUMERIC (Sum: {total:,.2f}, Avg: {avg:,.2f})\n"
            else:
                unique = df[col].nunique()
                report += f"  • {col}: TEXT ({unique} unique values)\n"
        
        # Totals section if we have numeric data
        if totals:
            report += f"\n**CALCULATED TOTALS:**\n"
            for col, total in totals.items():
                # Format based on value - whole numbers vs decimals
                if total == int(total):
                    report += f"  * **{col}**: {int(total):,}\n"
                else:
                    report += f"  * **{col}**: {total:,.2f}\n"
        
        # Grouping analysis - find text columns that might be used for grouping
        text_cols = [c for c in categorical_cols if c not in numeric_cols]
        if text_cols and numeric_cols:
            report += f"\n**GROUPING & ANALYSIS:**\n"
            for text_col in text_cols[:3]:  # Show top 3 text columns
                try:
                    # Group by text column and sum numeric columns
                    grouped = df.groupby(text_col)[numeric_cols].sum()
                    if len(grouped) <= 10:  # Only show if reasonable number of groups
                        report += f"\n  **By {text_col}:**\n"
                        for idx, (group_name, row) in enumerate(grouped.iterrows(), 1):
                            group_totals = " | ".join([f"{col}: {row[col]:,.2f}" if row[col] != int(row[col]) else f"{col}: {int(row[col]):,}" for col in numeric_cols])
                            report += f"    {idx}. {group_name}: {group_totals}\n"
                except Exception:
                    pass

        if text_cols:
            report += f"\n**CATEGORY COUNTS:**\n"
            for text_col in text_cols[:5]:
                try:
                    counts = df[text_col].fillna("Missing").replace("", "Missing").value_counts(dropna=False)
                    if len(counts) <= 20:
                        report += f"\n  **{text_col}:**\n"
                        for value, count in counts.items():
                            report += f"    - {value}: {int(count)}\n"
                except Exception:
                    pass
        
        # Show sample data
        report += f"\n**SAMPLE DATA (First 5 rows):**\n"
        for idx, row in df.head(5).iterrows():
            row_str = " | ".join([f"{col}: {row[col]}" for col in column_names])
            report += f"  {row_str}\n"

        report += "\n**RAW CSV DATA FOR CALCULATION:**\n```csv\n"
        report += csv_text.strip()
        report += "\n```\n"
        
        return _smart_truncate(report, "CSV")
    
    except ImportError:
        # pandas not available - fall back to basic CSV display
        logger.warning("pandas not available for CSV analysis - showing raw CSV")
        summary = f"📄 **CSV File**: {display_name}\n\n⚠️ Install pandas for intelligent CSV analysis:\n`pip install pandas`\n\n**Raw CSV:**\n{csv_text}"
        return _smart_truncate(summary, "CSV")
    except Exception as e:
        logger.error(f"Error analyzing CSV {display_name}: {e}", exc_info=True)
        # Fall back to raw display on error
        summary = f"📄 **CSV File**: {display_name}\n\n⚠️ Error during analysis: {str(e)}\n\n**Raw CSV:**\n{csv_text}"
        return _smart_truncate(summary, "CSV")


def process_attachment(attachment, corr_id: Optional[str] = None, user_id: Optional[str] = None,
                       raw_sink: Optional[dict] = None) -> str:
    """
    Process file attachment from Teams and extract text content.
    - Per Teams docs: downloadUrl in content is pre-authenticated by Teams
    - SharePoint/OneDrive contentUrl requires bot's Bearer token (app registration)
    - Returns extracted text for caller to cache and use in context
    - Implements size limits and truncation to prevent memory crashes
    
    Args:
        attachment: Teams attachment object
        corr_id: Correlation ID for logging
        user_id: User ID for logging purposes
        raw_sink: Optional dict to capture the raw downloaded bytes keyed by
            filename (used by the code interpreter for real file manipulation).
            Only populated on the direct-download path and capped to 25 MB.
    
    Returns:
        Extracted text or file info (truncated if necessary)
    """
    prefix = f"[{corr_id}] " if corr_id else ""
    
    try:
        display_name = getattr(attachment, "name", None) or "attachment"
        content_type = getattr(attachment, "content_type", None) or getattr(attachment, "contentType", "") or ""
        logger.info(f"{prefix}Processing: {display_name}")
        if content_type:
            logger.info(f"{prefix}Attachment content_type: {content_type}")
        
        # Per Teams documentation, files should have this structure:
        # - contentUrl: SharePoint/OneDrive URL (top-level)
        # - content.downloadUrl: Pre-authenticated download URL
        # - content.uniqueId: OneDrive item ID
        # - content.fileType: File extension
        # See: https://learn.microsoft.com/en-us/microsoftteams/platform/bots/how-to/bots-filesv4
        
        # Get file download URL from Teams attachment
        content = None
        download_url = None
        content_url = None
        content_info = None
        unique_id = None
        file_type = None
        
        # First: Parse content dict (Teams docs: contains downloadUrl and metadata)
        if hasattr(attachment, "content") and attachment.content:
            if isinstance(attachment.content, dict):
                content_info = attachment.content
            elif isinstance(attachment.content, str):
                s = attachment.content.strip()
                if s.startswith("{") or s.startswith("["):
                    try:
                        content_info = json.loads(s)
                        logger.info(f"{prefix}Parsed JSON string content into dict")
                    except Exception:
                        content_info = None
                else:
                    content_info = None
        
        # Extract fields from content dict
        if isinstance(content_info, dict):
            # Per Teams documentation: downloadUrl is the pre-authenticated URL
            if content_info.get("downloadUrl"):
                download_url = content_info.get("downloadUrl")
                logger.info(f"{prefix}Found pre-auth downloadUrl in content (Teams docs standard)")
            # Metadata extraction
            unique_id = content_info.get("uniqueId")
            file_type = content_info.get("fileType")
            if unique_id:
                logger.info(f"{prefix}Found uniqueId: {unique_id}")
            if file_type:
                logger.info(f"{prefix}Found fileType: {file_type}")
            # Improve display name if missing
            if display_name == "attachment":
                display_name = (
                    content_info.get("name")
                    or content_info.get("fileName")
                    or content_info.get("filename")
                    or display_name
                )
            # Fallback contentUrl from content dict
            if not download_url and content_info.get("contentUrl"):
                content_url = content_info.get("contentUrl")
                logger.info(f"{prefix}Found contentUrl in content dict (fallback)")
        
        # Second: Try top-level attributes (Teams docs: contentUrl is always present)
        if not content_url and hasattr(attachment, "contentUrl") and attachment.contentUrl:
            content_url = attachment.contentUrl
            logger.info(f"{prefix}Found contentUrl at top level (SharePoint/OneDrive)")
        elif not content_url and hasattr(attachment, "content_url") and attachment.content_url:
            content_url = attachment.content_url
            logger.info(f"{prefix}Found content_url at top level")
        
        # Third: Try direct downloadUrl attributes if not found in content
        if not download_url:
            if hasattr(attachment, "downloadUrl") and attachment.downloadUrl:
                download_url = attachment.downloadUrl
                logger.info(f"{prefix}Using attribute.downloadUrl")
            elif hasattr(attachment, "download_url") and attachment.download_url:
                download_url = attachment.download_url
                logger.info(f"{prefix}Using attribute.download_url")
            elif hasattr(attachment, "fileDownloadUrl") and attachment.fileDownloadUrl:
                download_url = attachment.fileDownloadUrl
                logger.info(f"{prefix}Using attribute.fileDownloadUrl")
        
        url_to_use = download_url or content_url
        if not url_to_use:
            # Debug: log all attachment attributes to identify the correct property
            logger.warning(f"{prefix}No URL found for attachment: {display_name}")
            logger.info(f"{prefix}  download_url={download_url}, content_url={content_url}")
            logger.info(f"{prefix}  content_info keys={list(content_info.keys()) if isinstance(content_info, dict) else 'N/A'}")
            
            # Get all non-private attributes
            all_attrs = [attr for attr in dir(attachment) if not attr.startswith('_')]
            logger.debug(f"{prefix}Attachment object attributes: {all_attrs}")
            
            # Log specific attribute values
            if hasattr(attachment, "content"):
                content_val = attachment.content
                logger.info(f"{prefix}  attachment.content type: {type(content_val).__name__}, len: {len(str(content_val)) if content_val else 0}")
                if isinstance(content_val, dict):
                    logger.info(f"{prefix}  Content dict keys: {list(content_val.keys())}")
                    for k, v in content_val.items():
                        if isinstance(v, str) and len(v) < 200:
                            logger.info(f"{prefix}    {k}={v}")
            
            if hasattr(attachment, "contentUrl"):
                logger.info(f"{prefix}  attachment.contentUrl: {attachment.contentUrl}")
            if hasattr(attachment, "content_url"):
                logger.info(f"{prefix}  attachment.content_url: {attachment.content_url}")
            if hasattr(attachment, "downloadUrl"):
                logger.info(f"{prefix}  attachment.downloadUrl: {attachment.downloadUrl}")
            if hasattr(attachment, "download_url"):
                logger.info(f"{prefix}  attachment.download_url: {attachment.download_url}")
            if hasattr(attachment, "fileDownloadUrl"):
                logger.info(f"{prefix}  attachment.fileDownloadUrl: {attachment.fileDownloadUrl}")
            
            # Log specific attribute values
            if hasattr(attachment, "content_type"):
                logger.info(f"{prefix}  attachment.content_type: {attachment.content_type}")
            if hasattr(attachment, "contentType"):
                logger.info(f"{prefix}  attachment.contentType: {attachment.contentType}")

            
            # Try alternative properties that might contain URLs
            url_props = [
                "url",
                "previewUrl",
                "preview_url",
                "fileUrl",
                "file_url",
                "thumbnailUrl",
                "thumbnail_url",
                "sourceUrl",
                "source_url",
                "fileDownloadUrl",
            ]
            
            for prop in url_props:
                if hasattr(attachment, prop):
                    val = getattr(attachment, prop)
                    if val:
                        logger.info(f"{prefix}Found URL in '{prop}': {val}")
                        url_to_use = val
                        break
            
            # If still no URL, check if content is a dict with any URL-like keys
            if not url_to_use and hasattr(attachment, "content") and isinstance(attachment.content, dict):
                for key, value in attachment.content.items():
                    if isinstance(value, str) and ("http://" in value or "https://" in value):
                        logger.info(f"{prefix}Found URL-like value in content['{key}']: {value}")
                        url_to_use = value
                        break
            
            if not url_to_use:
                    # Only skip if we have no file indicators at all
                    has_file_extension = display_name != "attachment" and any(
                        display_name.lower().endswith(ext) 
                        for ext in [".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt", ".txt", ".csv", ".json"]
                    )
                    
                    if not has_file_extension:
                        logger.warning(f"{prefix}No URL and no file extension found - skipping")
                        return "ℹ️ This looks like a link or card preview. Please upload the file directly so I can analyze it."
                    # Mobile-specific error message with guidance
                    return f"""❌ I detected file attachment '{display_name}' but couldn't access it.

**This often happens with mobile Teams app. Here's how to fix it:**

1. **Wait 10-30 seconds** after uploading, then send your message
   (Files need time to upload to OneDrive)

2. **Use the paperclip/attach button** instead of drag-and-drop

3. **Use Desktop or Web Teams** for best file attachment experience

4. **Check file size** (keep under 250 MB for best results)

The attachment payload did not include a usable download URL, which typically indicates a timing issue with mobile uploads."""
        
        # If SharePoint/OneDrive URL and no pre-auth downloadUrl, use robust Graph download+extract
        if _is_sharepoint_url(url_to_use) and not download_url:
            try:
                from knowledge_base import get_graph_token, download_and_extract_content
                token = get_graph_token()
                if not token:
                    logger.warning(f"{prefix}Graph token unavailable for SharePoint download")
                    return f"❌ Cannot access {display_name}: Graph API token unavailable."
                logger.info(f"{prefix}Downloading via Graph (shares/path) for SharePoint/OneDrive URL")
                return download_and_extract_content(url_to_use, token, display_name)
            except Exception as e:
                logger.error(f"{prefix}Graph download/extract error: {e}", exc_info=True)
                return f"❌ Error downloading {display_name} via Graph: {str(e)}"
        
        # Otherwise direct download and local extraction
        logger.info(f"{prefix}Downloading from: {url_to_use[:120]}...")
        import requests
        import time
        
        # SAFETY: Check file size first with HEAD request to avoid downloading huge files
        try:
            head_resp = requests.head(url_to_use, timeout=10, allow_redirects=True)
            content_length = head_resp.headers.get('Content-Length')
            if content_length:
                file_size_mb = int(content_length) / (1024 * 1024)
                logger.info(f"{prefix}File size: {file_size_mb:.1f} MB")
                
                if int(content_length) > (Config.MAX_FILE_SIZE_MB * 1024 * 1024):
                    logger.warning(f"{prefix}File too large: {file_size_mb:.1f} MB > {Config.MAX_FILE_SIZE_MB} MB")
                    return f"""❌ **File too large**: {display_name} ({file_size_mb:.1f} MB)

⚠️ Maximum file size: {Config.MAX_FILE_SIZE_MB} MB

**Suggestions:**
• Split the file into smaller parts
• Share specific sections instead of the full file
• Use a compressed format
• For very large datasets, consider uploading to SharePoint and asking specific questions"""
        except Exception as head_err:
            logger.debug(f"{prefix}HEAD request failed (continuing): {head_err}")
        
        # Mobile-friendly retry logic for attachment download
        max_retries = 3
        retry_delays = [2, 4, 8]  # Exponential backoff
        
        for attempt in range(max_retries):
            try:
                # Stream download with size limit for memory safety
                # REDUCED timeout: 15s to prevent long hangs that cause Teams timeout
                resp = requests.get(url_to_use, timeout=15, allow_redirects=True, stream=True)
                if resp.status_code == 200:
                    # Download with size limit to prevent memory crashes
                    content_chunks = []
                    total_size = 0
                    
                    for chunk in resp.iter_content(chunk_size=8192):
                        if chunk:
                            content_chunks.append(chunk)
                            total_size += len(chunk)
                            
                            # Safety: Stop if file exceeds limit during download
                            if total_size > (Config.MAX_FILE_SIZE_MB * 1024 * 1024):
                                logger.warning(f"{prefix}File exceeded size limit during download: {total_size / (1024*1024):.1f} MB")
                                return f"""❌ **File too large**: {display_name} (>{Config.MAX_FILE_SIZE_MB} MB)

⚠️ Download stopped to prevent memory issues.

**Please:**
• Upload a smaller file
• Share only relevant sections
• Use compressed formats"""
                    
                    content = b''.join(content_chunks)
                    logger.info(f"{prefix}Downloaded {len(content)} bytes ({len(content)/(1024*1024):.1f} MB) (attempt {attempt + 1})")
                    break
                elif resp.status_code == 403:
                    logger.error(f"{prefix}Download failed: HTTP {resp.status_code} (attempt {attempt + 1})")
                    if attempt < max_retries - 1:
                        logger.info(f"{prefix}Access denied, retrying in {retry_delays[attempt]}s (mobile timing issue?)")
                        time.sleep(retry_delays[attempt])
                        continue
                    return f"""❌ Access denied for {display_name} (HTTP 403).

**If using mobile Teams app:**
• Wait 30+ seconds after uploading, then try again
• Use the paperclip button (not drag-and-drop)  
• Switch to desktop/web Teams for better reliability

**Otherwise:**
• Check file permissions in OneDrive/SharePoint
• Make sure the file isn't restricted"""
                elif resp.status_code == 404:
                    if attempt < max_retries - 1:
                        logger.info(f"{prefix}File not found, retrying in {retry_delays[attempt]}s (upload may still be processing)")
                        time.sleep(retry_delays[attempt])
                        continue
                    return f"""❌ File not found: {display_name} (HTTP 404).

**If using mobile Teams app:**
• The file upload might not be complete yet
• Wait 30+ seconds after selecting the file, then try again
• Use desktop/web Teams for immediate file access"""
                else:
                    logger.error(f"{prefix}Download failed: HTTP {resp.status_code} (attempt {attempt + 1})")
                    if attempt < max_retries - 1:
                        time.sleep(retry_delays[attempt])
                        continue
                    return f"❌ Failed to download {display_name} (HTTP {resp.status_code})."
            except requests.exceptions.Timeout:
                logger.error(f"{prefix}Download timeout (>30s) (attempt {attempt + 1})")
                if attempt < max_retries - 1:
                    time.sleep(retry_delays[attempt])
                    continue
                return f"❌ Download timeout for {display_name}. File may be too large or network is slow."
            except Exception as e:
                logger.error(f"{prefix}Download exception: {e} (attempt {attempt + 1})", exc_info=True)
                if attempt < max_retries - 1:
                    time.sleep(retry_delays[attempt])
                    continue
                return f"❌ Failed to download {display_name}: {str(e)}"
        else:
            # All retries failed
            return f"""❌ Failed to download {display_name} after {max_retries} attempts.

**This commonly happens with mobile Teams:**
• Wait longer (60+ seconds) after uploading files
• Use desktop/web Teams instead of mobile app
• Check your network connection"""

        # Detect HTML viewer pages masquerading as files
        try:
            ctype = (resp.headers.get("Content-Type") or "").lower()
        except Exception:
            ctype = ""
        if ("text/html" in ctype) or (content[:20].lower().startswith(b"<html") or content[:40].lower().startswith(b"<!doctype html")):
            return f"❌ The downloaded content appears to be an HTML viewer page, not the raw file: {display_name}. If this is a SharePoint/OneDrive file, I can access it via Graph—please share the original file link."
        if not content:
            return f"📎 {display_name} (empty file)"
        
        # Capture raw bytes for the code interpreter (real xlsx/pdf/docx manipulation).
        # Capped to 25 MB to bound memory; larger files fall back to text-only.
        if raw_sink is not None:
            try:
                if len(content) <= 25 * 1024 * 1024:
                    raw_sink[display_name] = content
                else:
                    logger.info(f"{prefix}Raw bytes for '{display_name}' not captured ({len(content)} bytes > 25MB cap)")
            except Exception as _rb_err:
                logger.debug(f"{prefix}raw_sink capture failed: {_rb_err}")
        
        # Extract content from file bytes
        extracted_text = _extract_content(display_name, content)
        
        # Log successful extraction
        if extracted_text:
            logger.info(f"{prefix}Attachment {display_name} extracted - {len(extracted_text)} chars")
        
        return extracted_text
    
    except Exception as e:
        logger.error(f"{prefix}Error processing {getattr(attachment, 'name', 'file')}: {e}", exc_info=True)
        return f"❌ Error processing {getattr(attachment, 'name', 'attachment')}: {str(e)}"


def _extract_content(display_name: str, content: bytes) -> str:
    """Extract text content from file bytes based on file type.
    
    Implements smart truncation to prevent token limit crashes:
    - Extracts content normally
    - Truncates if exceeds MAX_EXTRACTED_CHARS
    - Adds clear indication when content is truncated
    """
    file_name = display_name.lower()
    
    # PDF with extensive analysis
    if file_name.endswith(".pdf"):
        if pypdf is None:
            return f"📄 PDF: {display_name}\n\n(Install pypdf to extract text.)"
        try:
            pdf_reader = pypdf.PdfReader(io.BytesIO(content))
            text_parts = []
            total_words = 0
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                try:
                    text = page.extract_text() or ""
                except Exception:
                    text = ""
                if text and text.strip():
                    cleaned = ' '.join(text.strip().split())
                    total_words += len(cleaned.split())
                    text_parts.append(f"Page {page_num}:\n{cleaned}")
            
            text = "\n\n".join(text_parts)
            
            # Generate extensive summary
            summary = f"📄 **PDF Document**: {display_name}\n"
            summary += f"**Pages:** {len(pdf_reader.pages)} | **Extracted Pages:** {len(text_parts)} | **Total Words:** {total_words:,}\n\n"
            summary += text
            return _smart_truncate(summary, "PDF")
        except Exception as e:
            return f"📄 PDF: {display_name}\n\n(Error: {str(e)})"
    
    # Word
    if file_name.endswith((".docx", ".doc")):
        if Document is None:
            return f"📝 Word: {display_name}\n\n(Install python-docx to extract text.)"
        
        # python-docx only works with .docx (Office Open XML format)
        if file_name.endswith(".doc") and not file_name.endswith(".docx"):
            # Try to detect if it's actually a .docx misnamed as .doc
            try:
                doc = Document(io.BytesIO(content))
                paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
                text = "\n\n".join(paragraphs)
                word_count = len(text.split())
                summary = f"📝 **Word Document**: {display_name}\n"
                summary += f"**Paragraphs:** {len(paragraphs)} | **Words:** {word_count:,}\n\n"
                summary += text
                return _smart_truncate(summary, "Word")
            except Exception as e:
                logger.warning(f"Error parsing {display_name} as DOCX: {e}")
                # This is a real legacy .doc file - try textract extraction
                legacy_text = _textract_legacy_office(content, "doc", display_name)
                if legacy_text:
                    word_count = len(legacy_text.split())
                    summary = f"📝 **Word Document (Legacy .doc)**: {display_name}\n"
                    summary += f"**Words:** {word_count:,} | **Size:** {len(content):,} bytes\n\n"
                    summary += legacy_text
                    return summary
                
                # Textract not available or failed
                if textract is None:
                    return (
                        f"📝 **Word (Legacy .doc)**: {display_name}\n\n"
                        f"⚠️ This is an old Word 97-2003 format file.\n"
                        f"Install textract to extract content: pip install textract\n"
                        f"Or convert to .docx format.\n"
                        f"File size: {len(content):,} bytes"
                    )
                return (
                    f"📝 **Word (Legacy .doc)**: {display_name}\n\n"
                    f"⚠️ Unable to extract text from this legacy .doc file.\n"
                    f"Please convert to .docx format for full extraction.\n"
                    f"File size: {len(content):,} bytes"
                )
        
        # .docx files with extensive summary
        try:
            doc = Document(io.BytesIO(content))
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
            text = "\n\n".join(paragraphs)
            
            # Generate extensive summary
            word_count = len(text.split())
            char_count = len(text)
            summary = f"📝 **Word Document**: {display_name}\n"
            summary += f"**Paragraphs:** {len(paragraphs)} | **Words:** {word_count:,} | **Characters:** {char_count:,}\n\n"
            summary += text
            return _smart_truncate(summary, "Word")
        except Exception as e:
            # Log detailed error for debugging
            error_msg = str(e)
            logger.error(f"Error extracting content from {display_name}: {error_msg}")
            # Check if this looks like a download/corruption issue
            if "not a" in error_msg.lower() or "is not" in error_msg.lower() or len(content) < 100:
                return (
                    f"📝 **Word Document**: {display_name}\n\n"
                    f"⚠️ The file appears to be corrupted or incomplete "
                    f"(size: {len(content):,} bytes, error: {error_msg[:80]})\n"
                    f"Please try re-uploading the file."
                )
            return f"📝 Word: {display_name}\n\n❌ Error: {error_msg[:200]}"
    
    # Excel - simple extraction preserving all data
    if file_name.endswith((".xlsx", ".xls")):
        # Legacy .xls
        if file_name.endswith(".xls") and not file_name.endswith(".xlsx"):
            if xlrd is None:
                return f"📊 Excel: {display_name}\n\n(Install xlrd==1.2.0 to extract legacy .xls content.)"
            try:
                wb = xlrd.open_workbook(file_contents=content)
                sheets_text = []
                
                for sheet in wb.sheets():
                    all_rows = []
                    
                    # Collect all rows
                    for rowx in range(sheet.nrows):
                        row = sheet.row_values(rowx)
                        formatted_row = [_format_cell_value(cell) for cell in row]
                        all_rows.append(formatted_row)
                    
                    if not all_rows:
                        continue
                    
                    # Build simple text output - tab-separated values
                    sheet_content = f"=== Sheet: {sheet.name} ===\n"
                    for row in all_rows:
                        sheet_content += "\t".join(row) + "\n"
                    sheets_text.append(sheet_content)
                
                result = f"📊 **Excel File**: {display_name}\n\n"
                result += "\n".join(sheets_text)
                return _smart_truncate(result, "Excel")
            except Exception as e:
                return f"📊 Excel: {display_name}\n\n(Error: {str(e)})"

        if load_workbook is None:
            return f"📊 Excel: {display_name}\n\n(Install openpyxl to extract content.)"
        try:
            import warnings
            with warnings.catch_warnings():
                warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl")
                wb = load_workbook(io.BytesIO(content), data_only=True)
            
            sheets_text = []
            
            for sheet in wb.worksheets:
                all_rows = []
                
                # Collect all rows
                for row in sheet.iter_rows(values_only=True):
                    formatted_row = [_format_cell_value(cell) for cell in row]
                    all_rows.append(formatted_row)
                
                if not all_rows:
                    continue
                
                # Build simple text output - tab-separated values
                sheet_content = f"=== Sheet: {sheet.title} ===\n"
                for row in all_rows:
                    sheet_content += "\t".join(row) + "\n"
                sheets_text.append(sheet_content)
            
            result = f"📊 **Excel File**: {display_name}\n\n"
            result += "\n".join(sheets_text)
            final_result = result if sheets_text else result + "[No data found]"
            return _smart_truncate(final_result, "Excel")
        except Exception as e:
            return f"📊 Excel: {display_name}\n\n(Error: {str(e)})"

    # PowerPoint
    if file_name.endswith((".pptx", ".ppt")):
        # python-pptx only supports .pptx (Office Open XML format), not legacy .ppt
        if file_name.endswith(".ppt") and not file_name.endswith(".pptx"):
            # Try to detect if it's actually a .pptx misnamed as .ppt
            try:
                # Check file signature (ZIP for .pptx, different for .ppt)
                if content[:4] == b'PK\x03\x04':
                    if Presentation is None:
                        return f"📽️ PowerPoint: {display_name}\n\n(Install python-pptx to extract content.)"
                    # This is a ZIP file (likely .pptx misnamed as .ppt)
                    logger.info(f"{display_name} appears to be .pptx format despite .ppt extension")
                    prs = Presentation(io.BytesIO(content))
                    slides_text = []
                    for idx, slide in enumerate(prs.slides, 1):
                        parts = []
                        for shape in slide.shapes:
                            try:
                                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                                    for p in shape.text_frame.paragraphs:
                                        txt = "".join(run.text for run in p.runs) if hasattr(p, "runs") else (p.text or "")
                                        if txt and txt.strip():
                                            parts.append(txt.strip())
                            except Exception as shape_error:
                                logger.debug(f"Error extracting from shape in slide {idx}: {shape_error}")
                                continue
                        if parts:
                            slides_text.append(f"Slide {idx}:\n" + "\n".join(parts))
                    text = "\n\n".join(slides_text)
                    text = text.strip()
                    if not text:
                        text = "[No extractable text found in slides]"
                    return f"📽️ **PowerPoint**: {display_name}\n\n{text}"
                else:
                    # This is a real legacy .ppt file
                    logger.warning(f"{display_name} is a legacy PowerPoint 97-2003 format")
                    legacy_text = _textract_legacy_office(content, "ppt", display_name)
                    if legacy_text:
                        # Estimate slide count from content structure (heuristic)
                        slide_markers = legacy_text.count('\n\n\n')  # Multiple line breaks often indicate slide boundaries
                        estimated_slides = max(1, slide_markers + 1)
                        word_count = len(legacy_text.split())
                        
                        summary = f"📽️ **PowerPoint (Legacy .ppt)**: {display_name}\n"
                        summary += f"**Estimated Slides:** {estimated_slides} | **Words:** {word_count:,} | **Size:** {len(content):,} bytes\n\n"
                        summary += legacy_text
                        return summary
                    
                    if textract is None:
                        return (
                            f"📽️ **PowerPoint (Legacy .ppt)**: {display_name}\n\n"
                            f"⚠️ Legacy .ppt requires extra tooling for text extraction.\n"
                            f"Install textract: pip install textract\n"
                            f"(May also require system dependencies like antiword, poppler-utils)\n"
                            f"File size: {len(content):,} bytes"
                        )
                    return (
                        f"📽️ **PowerPoint (Legacy .ppt)**: {display_name}\n\n"
                        f"⚠️ Unable to extract text from this legacy .ppt file.\n"
                        f"Textract is installed but extraction failed.\n"
                        f"File size: {len(content):,} bytes"
                    )
            except Exception as e:
                logger.warning(f"Error detecting format for {display_name}: {e}")
                legacy_text = _textract_legacy_office(content, "ppt", display_name)
                if legacy_text:
                    word_count = len(legacy_text.split())
                    summary = f"📽️ **PowerPoint (Legacy .ppt)**: {display_name}\n"
                    summary += f"**Words:** {word_count:,} | **Size:** {len(content):,} bytes\n\n"
                    summary += legacy_text
                    return summary
                
                if textract is None:
                    return (
                        f"📽️ **PowerPoint (Legacy .ppt)**: {display_name}\n\n"
                        f"⚠️ Legacy .ppt requires extra tooling. Install textract to enable extraction.\n"
                        f"File size: {len(content):,} bytes"
                    )
                return (
                    f"📽️ **PowerPoint (Legacy .ppt)**: {display_name}\n\n"
                    f"⚠️ Unable to extract text from this legacy .ppt file.\n"
                    f"File size: {len(content):,} bytes"
                )
        if Presentation is None:
            return f"📽️ PowerPoint: {display_name}\n\n(Install python-pptx to extract content.)"
        
        # .pptx files - standard extraction
        try:
            # Attempt to load presentation
            try:
                prs = Presentation(io.BytesIO(content))
            except (AttributeError, TypeError) as e:
                # Known python-pptx issue with malformed relationships (rId errors)
                if "rId" in str(e) or "'list' object has no attribute" in str(e):
                    logger.warning(f"PowerPoint extraction error for {display_name}: {e} - file may have malformed relationships")
                    return f"📽️ PowerPoint: {display_name}\n\n⚠️ This PowerPoint file has a structural issue (malformed relationships) that prevents text extraction. The file may still be viewable in PowerPoint."
                raise
            
            slides_text = []
            total_shapes = 0
            total_text_items = 0
            
            for idx, slide in enumerate(prs.slides, 1):
                parts = []
                for shape in slide.shapes:
                    total_shapes += 1
                    try:
                        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                            for p in shape.text_frame.paragraphs:
                                txt = "".join(run.text for run in p.runs) if hasattr(p, "runs") else (p.text or "")
                                if txt and txt.strip():
                                    parts.append(txt.strip())
                                    total_text_items += 1
                    except Exception as shape_error:
                        # Skip problematic shapes but log them
                        logger.debug(f"Error extracting from shape in slide {idx}: {shape_error}")
                        continue
                if parts:
                    slides_text.append(f"**Slide {idx}** ({len(parts)} text items):\n" + "\n".join(parts))
            
            text = "\n\n".join(slides_text)
            text = text.strip()
            
            # Generate extensive summary
            summary = f"📽️ **PowerPoint**: {display_name}\n"
            summary += f"**Total Slides:** {len(prs.slides)} | **Shapes:** {total_shapes} | **Text Items:** {total_text_items}\n\n"
            
            if not text:
                summary += "[No extractable text found in slides]"
            else:
                summary += text
            
            return _smart_truncate(summary, "PowerPoint")
        except Exception as e:
            logger.error(f"PowerPoint extraction error for {display_name}: {e}", exc_info=True)
            error_msg = str(e)
            # Check if this is a file format error
            if "not a" in error_msg.lower() or "is not" in error_msg.lower() or "invalid" in error_msg.lower():
                return (
                    f"📽️ PowerPoint: {display_name}\n\n"
                    f"⚠️ Unable to extract text from this file. It may be:\n"
                    f"• A legacy .ppt format (PowerPoint 97-2003)\n"
                    f"• Corrupted or incomplete\n"
                    f"• Password protected\n\n"
                    f"Try opening in PowerPoint or converting to .pptx format.\n"
                    f"Error: {error_msg[:150]}"
                )
            return f"📽️ PowerPoint: {display_name}\n\n(Error: {error_msg[:200]})"
    
    # Text files - simple extraction
    if file_name.endswith((".txt", ".md", ".json", ".xml", ".csv", ".log")):
        try:
            text = content.decode("utf-8", errors="ignore").strip()
            
            # CSV - intelligent analysis with totals, summaries, grouping
            if file_name.endswith(".csv"):
                if not text:
                    return f"📄 **CSV File**: {display_name}\n\n[Empty file]"
                return _analyze_csv_content(text, display_name)
            
            # JSON with structure info
            elif file_name.endswith(".json"):
                label = "JSON Data"
                try:
                    import json as json_lib
                    data = json_lib.loads(text)
                    if isinstance(data, dict):
                        keys_info = f" | **Keys:** {len(data)} ({', '.join(list(data.keys())[:10])})"
                    elif isinstance(data, list):
                        keys_info = f" | **Items:** {len(data)}"
                    else:
                        keys_info = ""
                    summary = f"📄 **{label}**: {display_name}{keys_info}\n\n{text}"
                    return _smart_truncate(summary, "JSON")
                except:
                    pass
            
            # XML with size info
            elif file_name.endswith(".xml"):
                label = "XML Data"
                summary = f"📄 **{label}**: {display_name} | **Size:** {len(text):,} chars\n\n{text}"
                return _smart_truncate(summary, "XML")
            
            # Default text file
            label = "Text File"
            summary = f"📄 **{label}**: {display_name}\n\n{text}"
            return _smart_truncate(summary, "Text")
        except Exception as e:
            return f"📄 Text: {display_name}\n\n(Error: {str(e)})"
    
    # Images
    if file_name.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp")):
        if Image is None:
            return f"🖼️ Image: {display_name}\n\n(Install pillow to inspect image.)"
        try:
            img = Image.open(io.BytesIO(content))
            img_info = f"🖼️ **Image**: {display_name} ({img.width}x{img.height}px, {img.format})"
            
            # Analyze image using Azure OpenAI vision
            analysis = _analyze_image_vision(content, display_name)
            if analysis:
                return f"{img_info}\n\n**Analysis:**\n{analysis}"
            return img_info
        except Exception as e:
            return f"🖼️ Image: {display_name}\n\n(Error: {str(e)})"


def _analyze_image_vision(image_bytes: bytes, display_name: str) -> str:
    """
    Analyze image using Azure OpenAI vision API.
    Returns description of image contents.
    """
    try:
        import os
        from openai import AzureOpenAI
        
        client = AzureOpenAI(
            api_key=os.getenv("AZURE_OPENAI_API_KEY"),
            api_version="2024-02-01",
            azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT")
        )
        
        # Encode image to base64
        image_b64 = base64.b64encode(image_bytes).decode("utf-8")
        
        # Determine media type
        file_name_lower = display_name.lower()
        if file_name_lower.endswith(".png"):
            media_type = "image/png"
        elif file_name_lower.endswith(".gif"):
            media_type = "image/gif"
        elif file_name_lower.endswith(".bmp"):
            media_type = "image/bmp"
        else:  # .jpg, .jpeg, or default
            media_type = "image/jpeg"
        
        # Call vision API with text extraction focus
        response = client.chat.completions.create(
            model=os.getenv("AZURE_OPENAI_MODEL_DEPLOYMENT_NAME", "gpt-4o"),
            max_tokens=1000,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{media_type};base64,{image_b64}"
                            }
                        },
                        {
                            "type": "text",
                            "text": """Extract and transcribe ALL visible text, labels, numbers, titles, and data from this image. Include:
- Chart/diagram titles and headings
- All axis labels and values
- Data point labels and numeric values
- Legend items and keys
- Any annotations or callouts
- Table data (if present)
- All readable text in any part of the image

Format as structured text. Be precise with numbers and labels. Do not summarize or guess—extract exactly what is visible."""
                        }
                    ]
                }
            ]
        )
        
        if response.choices and len(response.choices) > 0:
            analysis = response.choices[0].message.content.strip()
            logger.info(f"Image analysis successful for {display_name}")
            return analysis
        return None
    except Exception as e:
        logger.warning(f"Image vision analysis failed for {display_name}: {e}")
        return None


def search_local_files(query: str, file_list: list) -> str:
    """
    Simple keyword search in previously uploaded files (stored in conversation memory).
    No external storage access needed.
    
    Args:
        query: Search query
        file_list: List of previously processed files in conversation
    
    Returns:
        Search results
    """
    if not file_list:
        return "\n\n❌ No files uploaded yet. Please upload files first."
    
    query_lower = query.lower() if query else ""
    results = []
    
    for file_info in file_list:
        file_name = file_info.get("name") or ""
        file_content = file_info.get("content") or ""
        
        # Skip if both are empty
        if not file_name and not file_content:
            continue
        
        # Safe search with None checks
        if query_lower and (query_lower in file_name.lower() or query_lower in file_content.lower()):
            results.append({
                "name": file_name or "Unnamed file",
                "snippet": file_content[:200] + "..." if len(file_content) > 200 else file_content
            })
    
    if not results:
        return f"\n\n❌ No files found matching '{query}' in uploaded files."
    
    output = f"\n\n🔍 **Search Results** (found {len(results)} file(s)):\n"
    for r in results[:5]:
        output += f"\n• **{r['name']}**\n  {r['snippet']}\n"
    
    return output


def aggregate_tabular_files(extracted_contents: list[tuple[str, str]]) -> str:
    """
    Compare and aggregate multiple CSV/Excel files.
    Works with any tabular data - timesheets, inventory, sales, etc.
    
    Args:
        extracted_contents: List of (filename, extracted_content) tuples
    
    Returns:
        Combined comparison analysis
    """
    if not extracted_contents:
        return ""
    
    # Only aggregate tabular files (CSV, Excel)
    tabular_files = []
    for filename, content in extracted_contents:
        fname_lower = filename.lower()
        if fname_lower.endswith(('.csv', '.xlsx', '.xls')):
            tabular_files.append((filename, content))
    
    # If less than 2 tabular files, no comparison needed
    if len(tabular_files) < 2:
        return ""
    
    # Parse each file
    file_data = []
    all_columns = set()
    
    for filename, content in tabular_files:
        file_info = {
            "name": filename,
            "rows": 0,
            "columns": [],
            "column_stats": {},
            "persons": {},
            "sample_data": []
        }
        
        lines = content.split('\n')
        in_stats_section = False
        in_person_section = False
        in_data_section = False
        
        for line in lines:
            line_stripped = line.strip()
            
            # Count rows from metadata like "**Rows:** 237"
            rows_match = re.search(r'\*\*Rows[:\*]*\s*(\d+)', line, re.IGNORECASE)
            if rows_match:
                file_info["rows"] = int(rows_match.group(1))
            
            # Extract column headers from "**Column Headers:** col1, col2, ..."
            if 'Column Headers' in line or 'Columns:' in line:
                cols_match = re.search(r'[:\*]+\s*(.+)$', line)
                if cols_match:
                    cols = [c.strip() for c in cols_match.group(1).split(',')]
                    file_info["columns"] = cols
                    all_columns.update(cols)
            
            # Detect sections
            if '📊' in line and ('Statistics' in line or 'COLUMN' in line):
                in_stats_section = True
                in_person_section = False
                in_data_section = False
                continue
            
            if '👥' in line or 'By Person' in line or 'HOURS BY PERSON' in line:
                in_person_section = True
                in_stats_section = False
                in_data_section = False
                continue
            
            if 'FULL DATA' in line or 'Data Rows' in line:
                in_data_section = True
                in_stats_section = False
                in_person_section = False
                continue
            
            # End sections
            if not line_stripped or line_stripped.startswith('---') or line_stripped.startswith('==='):
                if not in_data_section:
                    in_stats_section = False
                    in_person_section = False
                continue
            
            # Parse column statistics
            if in_stats_section and ('•' in line_stripped or line_stripped.startswith('-')):
                try:
                    entry = line_stripped.lstrip('•-').strip().replace('**', '')
                    if ':' in entry and 'SUM=' in entry.upper():
                        col_name, stats_part = entry.split(':', 1)
                        col_name = col_name.strip()
                        
                        sum_match = re.search(r'SUM\s*=\s*([\d,.-]+)', stats_part, re.IGNORECASE)
                        count_match = re.search(r'COUNT\s*=\s*(\d+)', stats_part, re.IGNORECASE)
                        avg_match = re.search(r'AVG\s*=\s*([\d,.-]+)', stats_part, re.IGNORECASE)
                        min_match = re.search(r'MIN\s*=\s*([\d,.-]+)', stats_part, re.IGNORECASE)
                        max_match = re.search(r'MAX\s*=\s*([\d,.-]+)', stats_part, re.IGNORECASE)
                        
                        file_info["column_stats"][col_name] = {
                            "sum": float(sum_match.group(1).replace(',', '')) if sum_match else 0,
                            "count": int(count_match.group(1)) if count_match else 0,
                            "avg": float(avg_match.group(1).replace(',', '')) if avg_match else 0,
                            "min": float(min_match.group(1).replace(',', '')) if min_match else 0,
                            "max": float(max_match.group(1).replace(',', '')) if max_match else 0,
                        }
                        all_columns.add(col_name)
                except Exception:
                    pass
            
            # Parse person breakdown
            if in_person_section and ('•' in line_stripped or line_stripped.startswith('-')):
                try:
                    entry = line_stripped.lstrip('•-').strip().replace('**', '')
                    if ':' in entry:
                        name_part, rest = entry.split(':', 1)
                        person = name_part.strip()
                        
                        # Extract numeric value (hours, amount, count, etc.)
                        value_match = re.search(r'([\d,]+\.?\d*)', rest)
                        entries_match = re.search(r'\((\d+)\s*entr', rest.lower())
                        
                        if person and value_match:
                            file_info["persons"][person] = {
                                "value": float(value_match.group(1).replace(',', '')),
                                "entries": int(entries_match.group(1)) if entries_match else 1
                            }
                except Exception:
                    pass
            
            # Capture sample data rows (first 3)
            if in_data_section and line_stripped.startswith('Row'):
                if len(file_info["sample_data"]) < 3:
                    file_info["sample_data"].append(line_stripped[:200])
        
        file_data.append(file_info)
    
    # Build comparison output
    output = []
    output.append("=" * 60)
    output.append("📊 **MULTI-FILE COMPARISON**")
    output.append("=" * 60)
    output.append("")
    
    # Per-file summary
    output.append(f"**📁 FILES ({len(tabular_files)}):**")
    output.append("")
    
    for fd in file_data:
        output.append(f"**{fd['name']}**")
        if fd['rows']:
            output.append(f"  • Rows: {fd['rows']}")
        if fd['columns']:
            output.append(f"  • Columns: {', '.join(fd['columns'][:8])}" + (f" (+{len(fd['columns'])-8} more)" if len(fd['columns']) > 8 else ""))
        
        # Show column statistics
        if fd['column_stats']:
            output.append(f"  • Statistics:")
            for col, stats in list(fd['column_stats'].items())[:5]:  # Top 5 columns
                output.append(f"      - {col}: SUM={stats['sum']:.2f}, AVG={stats['avg']:.2f}, COUNT={stats['count']}")
        
        # Show top people/items if available
        if fd['persons']:
            sorted_persons = sorted(fd['persons'].items(), key=lambda x: x[1]["value"], reverse=True)[:5]
            output.append(f"  • Top Entries:")
            for person, data in sorted_persons:
                output.append(f"      - {person}: {data['value']:.2f} ({data['entries']} entries)")
        
        output.append("")
    
    # Cross-file comparison for common numeric columns
    common_stats_cols = set()
    for fd in file_data:
        common_stats_cols.update(fd['column_stats'].keys())
    
    if common_stats_cols:
        output.append("**📊 COLUMN COMPARISON ACROSS FILES:**")
        output.append("")
        
        for col in common_stats_cols:
            col_values = []
            for fd in file_data:
                if col in fd['column_stats']:
                    stats = fd['column_stats'][col]
                    col_values.append({
                        "file": fd['name'][:30],
                        "sum": stats['sum'],
                        "count": stats['count'],
                        "avg": stats['avg']
                    })
            
            if len(col_values) >= 2:
                output.append(f"**{col}:**")
                total_sum = sum(cv['sum'] for cv in col_values)
                total_count = sum(cv['count'] for cv in col_values)
                for cv in col_values:
                    pct = (cv['sum'] / total_sum * 100) if total_sum > 0 else 0
                    output.append(f"  • {cv['file']}: SUM={cv['sum']:.2f} ({pct:.1f}%), COUNT={cv['count']}")
                output.append(f"  • **TOTAL: {total_sum:.2f}** ({total_count} records)")
                output.append("")
    
    # Combined person/entity breakdown across all files
    combined_persons = {}
    for fd in file_data:
        for person, data in fd['persons'].items():
            if person not in combined_persons:
                combined_persons[person] = {"value": 0, "entries": 0, "files": []}
            combined_persons[person]["value"] += data["value"]
            combined_persons[person]["entries"] += data["entries"]
            combined_persons[person]["files"].append(fd['name'][:20])
    
    if combined_persons:
        output.append("**👥 COMBINED BY NAME/ENTITY (All Files):**")
        sorted_combined = sorted(combined_persons.items(), key=lambda x: x[1]["value"], reverse=True)
        total_value = sum(p[1]["value"] for p in sorted_combined)
        
        for person, data in sorted_combined[:15]:  # Top 15
            pct = (data["value"] / total_value * 100) if total_value > 0 else 0
            files_count = len(data["files"])
            files_note = f" (in {files_count} files)" if files_count > 1 else ""
            output.append(f"  • **{person}**: {data['value']:.2f} ({data['entries']} entries) - {pct:.1f}%{files_note}")
        
        if len(sorted_combined) > 15:
            output.append(f"  • ... and {len(sorted_combined) - 15} more")
        
        output.append("")
        output.append(f"**GRAND TOTAL: {total_value:.2f}**")
    
    output.append("")
    output.append("=" * 60)
    output.append("")
    
    return "\n".join(output)
