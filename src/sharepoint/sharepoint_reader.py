from __future__ import annotations

import csv
import io
import logging
import os
import time

import requests

from config import Config
from sharepoint.graph_client import download_file_bytes, list_configured_sharepoint_drives, list_drive_items

logger = logging.getLogger(__name__)

# Document Intelligence (prebuilt-read) handles these server-side, including OCR for
# scanned/image PDFs and photographs. Office formats are extracted locally (faster and
# preserves structure), text/csv are read directly.
SUPPORTED_EXTENSIONS = {
    # Office formats — modern (Open XML)
    ".pdf", ".docx", ".xlsx", ".pptx",
    # Office formats — legacy
    ".doc", ".xls", ".ppt",
    # Text / markup / data
    ".txt", ".md", ".csv", ".tsv", ".json", ".xml", ".log", ".rtf",
    ".html", ".htm", ".yaml", ".yml",
    # Images (OCR + table extraction via Document Intelligence)
    ".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp",
}
# Formats Document Intelligence is the primary extractor for (OCR + tables).
_DI_PRIMARY_EXTENSIONS = {
    ".pdf", ".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp",
}
_DI_API_VERSION = "2024-11-30"


def _extract_with_document_intelligence(file_bytes: bytes, file_name: str) -> str:
    """Extract text via Azure Document Intelligence (prebuilt-read) REST API.

    Returns the analyzed text, or "" on any failure so the caller can fall back to
    local extraction. Uses the REST API directly (no extra SDK dependency) so it works
    unchanged inside the indexer subprocess.
    """
    endpoint = (Config.AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT or "").rstrip("/")
    key = Config.AZURE_DOCUMENT_INTELLIGENCE_KEY
    if not endpoint or not key:
        return ""
    model = Config.DOCUMENT_INTELLIGENCE_MODEL or "prebuilt-read"
    analyze_url = (
        f"{endpoint}/documentintelligence/documentModels/{model}:analyze"
        f"?api-version={_DI_API_VERSION}"
    )
    headers = {
        "Ocp-Apim-Subscription-Key": key,
        "Content-Type": "application/octet-stream",
    }
    deadline = time.monotonic() + Config.DOCUMENT_INTELLIGENCE_TIMEOUT
    try:
        resp = requests.post(analyze_url, headers=headers, data=file_bytes, timeout=Config.DOCUMENT_INTELLIGENCE_TIMEOUT)
        resp.raise_for_status()
        operation_url = resp.headers.get("Operation-Location")
        if not operation_url:
            logger.warning("Document Intelligence: no Operation-Location for %s", file_name)
            return ""
        poll_headers = {"Ocp-Apim-Subscription-Key": key}
        while time.monotonic() < deadline:
            time.sleep(1.5)
            poll = requests.get(operation_url, headers=poll_headers, timeout=Config.DOCUMENT_INTELLIGENCE_TIMEOUT)
            poll.raise_for_status()
            body = poll.json()
            status = (body.get("status") or "").lower()
            if status == "succeeded":
                content = (body.get("analyzeResult") or {}).get("content", "") or ""
                logger.info(
                    "Document Intelligence OK | file=%s | model=%s | chars=%s",
                    file_name, model, len(content),
                )
                return content
            if status in ("failed", "canceled"):
                logger.warning("Document Intelligence %s for %s: %s", status, file_name, body.get("error"))
                return ""
        logger.warning("Document Intelligence timed out for %s after %ss", file_name, Config.DOCUMENT_INTELLIGENCE_TIMEOUT)
        return ""
    except Exception as exc:
        logger.warning("Document Intelligence failed for %s (%s) — falling back to local extraction", file_name, type(exc).__name__)
        return ""


def _extract_pdf_local(file_bytes: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    return "\n".join(page.extract_text() or "" for page in reader.pages)



def list_sharepoint_documents(max_items: int | None = None, max_depth: int = 8) -> list[dict]:
    documents: list[dict] = []
    for drive in list_configured_sharepoint_drives():
        remaining = None if max_items is None else max(max_items - len(documents), 0)
        if remaining == 0:
            break
        documents.extend(list_drive_items(drive["site_id"], drive["drive_id"], max_items=remaining, max_depth=max_depth))
    return documents


def download_sharepoint_document(item: dict) -> bytes:
    return download_file_bytes(item["drive_id"], item["id"])


def extract_document_text(file_name: str, file_bytes: bytes) -> str:
    ext = os.path.splitext(file_name or "")[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        logger.info("Skipping unsupported SharePoint file type: %s", file_name)
        return ""

    # Document Intelligence is the primary extractor for PDFs and images (server-side
    # OCR + table extraction, no local GIL pressure). Fall back to local pypdf for PDFs
    # when DI is disabled or returns nothing; images have no local fallback.
    if ext in _DI_PRIMARY_EXTENSIONS and Config.ENABLE_DOCUMENT_INTELLIGENCE:
        di_text = _extract_with_document_intelligence(file_bytes, file_name)
        if di_text.strip():
            return di_text
        if ext == ".pdf":
            logger.info("Document Intelligence returned no text for %s — using local pypdf", file_name)
        else:
            return ""  # image OCR only available via DI

    if ext == ".txt":
        return file_bytes.decode("utf-8", errors="ignore")
    if ext == ".csv":
        text = file_bytes.decode("utf-8", errors="ignore")
        return "\n".join(", ".join(row) for row in csv.reader(io.StringIO(text)))
    if ext == ".pdf":
        return _extract_pdf_local(file_bytes)
    if ext == ".docx":
        from docx import Document

        doc = Document(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if ext == ".xlsx":
        import openpyxl

        workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        lines: list[str] = []
        for sheet in workbook.worksheets:
            lines.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [str(cell) for cell in row if cell is not None]
                if values:
                    lines.append(", ".join(values))
        return "\n".join(lines)
    if ext == ".pptx":
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(file_bytes))
        lines: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text)
        return "\n".join(lines)

    # ---- Legacy Office formats ------------------------------------------------
    if ext == ".doc":
        # python-docx can read many .doc files saved as Open XML under a .doc name.
        # Real legacy binary .doc (pre-2007) will raise; we return empty so the
        # caller skips the document rather than crashing the indexing run.
        try:
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            return ""

    if ext == ".xls":
        try:
            import xlrd
            wb = xlrd.open_workbook(file_contents=file_bytes)
            lines: list[str] = []
            for sheet in wb.sheets():
                lines.append(f"Sheet: {sheet.name}")
                for rowx in range(sheet.nrows):
                    values = [str(v) for v in sheet.row_values(rowx) if str(v).strip()]
                    if values:
                        lines.append(", ".join(values))
            return "\n".join(lines)
        except Exception:
            return ""

    if ext == ".ppt":
        # .pptx disguised as .ppt (ZIP magic bytes)
        try:
            if file_bytes[:4] == b"PK\x03\x04":
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_bytes))
                lines: list[str] = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            lines.append(shape.text)
                return "\n".join(lines)
        except Exception:
            pass
        return ""

    # ---- Text / markup / data formats ----------------------------------------
    if ext == ".md":
        return file_bytes.decode("utf-8", errors="ignore")

    if ext in (".json", ".xml", ".log"):
        return file_bytes.decode("utf-8", errors="ignore")

    if ext in (".html", ".htm"):
        import re as _re
        from html.parser import HTMLParser

        class _Ex(HTMLParser):
            _SKIP = frozenset({"script", "style", "head"})

            def __init__(self):
                super().__init__()
                self._p: list[str] = []
                self._depth = 0

            def handle_starttag(self, tag, attrs):
                if tag.lower() in self._SKIP:
                    self._depth += 1

            def handle_endtag(self, tag):
                if tag.lower() in self._SKIP:
                    self._depth = max(0, self._depth - 1)

            def handle_data(self, data):
                if not self._depth:
                    self._p.append(data)

            def text(self):
                return _re.sub(r"\s+", " ", "".join(self._p)).strip()

        p = _Ex()
        p.feed(file_bytes.decode("utf-8", errors="ignore"))
        return p.text()

    if ext == ".rtf":
        import re as _re
        raw = file_bytes.decode("latin-1", errors="ignore")
        text = _re.sub(r"\\[a-z]+\-?\d*\s?", " ", raw)
        text = _re.sub(r"\\.", " ", text)
        text = _re.sub(r"[{}]", " ", text)
        return _re.sub(r"\s+", " ", text).strip()

    if ext == ".tsv":
        text = file_bytes.decode("utf-8", errors="ignore")
        return "\n".join(", ".join(row) for row in csv.reader(io.StringIO(text), delimiter="\t"))

    if ext in (".yaml", ".yml"):
        return file_bytes.decode("utf-8", errors="ignore")

    return ""
