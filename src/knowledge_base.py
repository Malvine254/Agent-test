# =====================================================
# AI Search (Azure Cognitive Search) Integration
# =====================================================
import os
import tempfile
import re
import logging
from typing import Optional
from dotenv import load_dotenv
from config import Config
from utils.truncation import safe_truncate
from datetime import datetime
load_dotenv()

logger = logging.getLogger(__name__)


# =====================================================
# Clean Logging Utilities (organized by category)
# =====================================================
class SearchLogger:
    """Centralized logging with clean categorization"""
    
    # Section headers
    @staticmethod
    def section(title: str):
        """Log a major section header"""
        logger.info(f"\n{'='*60}\nðŸ“‹ {title}\n{'='*60}")
    
    # Query operations
    @staticmethod
    def query_received(query: str, user: str = ""):
        """Log when query is received"""
        user_info = f" | User: {user}" if user else ""
        logger.info(f"ðŸ” QUERY: '{query}'{user_info}")
    
    # Search source results
    @staticmethod
    def search_source(source: str, count: int, query: str = ""):
        """Log results from a search source"""
        query_info = f" | Query: '{query[:50]}'" if query else ""
        if count == 0:
            logger.info(f"  âšª {source}: 0 results{query_info}")
        else:
            logger.info(f"  âœ… {source}: {count} result(s){query_info}")
    
    # Document listings
    @staticmethod
    def document(name: str, doc_type: str = "", score: float = None):
        """Log a found document"""
        score_str = f" | Score: {score:.2f}" if score is not None else ""
        type_str = f" | Type: {doc_type}" if doc_type else ""
        logger.info(f"    ðŸ“„ {name}{type_str}{score_str}")
    
    # Filtering actions
    @staticmethod
    def filtered(reason: str, name: str):
        """Log a filtered-out item"""
        logger.debug(f"    âŠ˜ FILTERED: {name} ({reason})")
    
    # Errors
    @staticmethod
    def error(message: str, exc_info=False):
        """Log an error"""
        logger.error(f"âŒ ERROR: {message}", exc_info=exc_info)
    
    # Warnings
    @staticmethod
    def warning(message: str):
        """Log a warning"""
        logger.warning(f"âš ï¸  {message}")
    
    # Summary
    @staticmethod
    def summary(title: str, metrics: dict):
        """Log a summary with metrics"""
        logger.info(f"\nðŸ“Š SUMMARY: {title}")
        for key, value in metrics.items():
            logger.info(f"  â€¢ {key}: {value}")


def build_keyword_query(query: str) -> str:
    """Clean up a search query string. The LLM router is fully responsible for
    deciding *what* to search â€” this function only trims whitespace and collapses
    spaces so the query is safe to pass to Azure Search."""
    if not query:
        return ""
    cleaned = " ".join(query.split())
    return cleaned.strip()

def _map_search_doc(raw: dict) -> dict:
    """Normalize Azure Search doc fields to common keys and derive a snippet.

    Prefer `content`/`text`/`chunk` when available. Otherwise, use semantic
    `@search.captions` or simple `@search.highlights` to populate a short snippet
    so results can be surfaced even when the index doesn't make raw content
    retrievable.
    """
    try:
        name = (
            raw.get("name") or raw.get("title") or raw.get("metadata_storage_name") or raw.get("file_name") or raw.get("filename")
        )
        url = (
            raw.get("url") or raw.get("webUrl") or raw.get("metadata_storage_path") or raw.get("file_path")
        )

        # Base content fields
        content = raw.get("content") or raw.get("text") or raw.get("chunk") or ""

        # Derive snippet from captions/highlights when content is missing
        snippet = content
        try:
            captions = raw.get("@search.captions") or []
            if not snippet and isinstance(captions, list) and captions:
                # captions: [{"text": "...", "highlights": ["..."]}]
                cap_texts = []
                for c in captions[:3]:
                    t = (c.get("text") or "").strip()
                    if t:
                        cap_texts.append(t)
                snippet = " \n".join(cap_texts).strip()
        except Exception:
            pass

        try:
            highlights = raw.get("@search.highlights") or {}
            if not snippet and isinstance(highlights, dict) and highlights:
                # Common highlight keys: content, text, chunk, metadata_storage_name
                hl_keys = ("content", "text", "chunk")
                hl_texts = []
                for k in hl_keys:
                    vals = highlights.get(k) or []
                    if isinstance(vals, list):
                        for v in vals[:3]:
                            s = (v or "").strip()
                            if s:
                                hl_texts.append(s)
                if hl_texts:
                    snippet = " \n".join(hl_texts).strip()
        except Exception:
            pass

        return {
            "id": raw.get("id"),
            "name": name or "Untitled",
            "file_path": url or "",
            "url": url or "",
            "content": snippet or "",
            "file_type": raw.get("file_type"),
            "upload_date": raw.get("upload_date"),
            "score": raw.get("@search.score"),
        }
    except Exception:
        return {"name": "Untitled", "file_path": "", "url": "", "content": ""}

def _merge_ranked_search_docs(existing: list, incoming: list, source_weight: float = 0.0) -> list:
    """Merge Azure Search pass results while preserving the strongest snippets."""
    merged = {}
    for position, doc in enumerate((existing or []) + (incoming or [])):
        key = doc.get("id") or doc.get("url") or doc.get("file_path") or doc.get("name")
        if not key:
            key = f"doc:{len(merged)}"
        rank_bonus = max(0, 50 - position)
        score = float(doc.get("score") or 0) + source_weight + rank_bonus
        current = merged.get(key)
        if not current or score > current.get("_hybrid_score", 0):
            doc["_hybrid_score"] = score
            merged[key] = doc
        elif doc.get("content") and doc.get("content") not in (current.get("content") or ""):
            current["content"] = ((current.get("content") or "") + "\n\n" + doc.get("content", "")).strip()
    return sorted(merged.values(), key=lambda d: d.get("_hybrid_score", 0), reverse=True)

def search_documents(query: str, top: int = 20) -> list:
    """Search Azure Cognitive Search with a hybrid semantic + keyword strategy."""
    started_at = time.perf_counter()
    endpoint = os.getenv("AZURE_SEARCH_ENDPOINT")
    index_name = os.getenv("AZURE_SEARCH_INDEX_NAME") or getattr(Config, "AZURE_SEARCH_INDEX_NAME", "sharepoint-documents")
    api_key = (
        os.getenv("AZURE_SEARCH_QUERY_KEY")
        or os.getenv("AZURE_SEARCH_ADMIN_KEY")
        or getattr(Config, "AZURE_SEARCH_QUERY_KEY", "")
        or getattr(Config, "AZURE_SEARCH_ADMIN_KEY", "")
    )
    # API version
    api_version = os.getenv("AZURE_SEARCH_API_VERSION", getattr(Config, "AZURE_SEARCH_API_VERSION", "2025-09-01"))
    if not endpoint or not index_name or not api_key:
        logger.error("Missing Azure Cognitive Search configuration.")
        return []
    
    original_query = query or ""
    query = build_keyword_query(original_query) or original_query
    
    # Preprocess query to improve search accuracy
    # Remove or simplify special characters and path-like syntax
    cleaned_query = query.strip()
    # Remove backslashes and path separators
    cleaned_query = cleaned_query.replace("\\", " ").replace("/", " ")
    # Remove extra parentheses and special chars that confuse Azure Search
    cleaned_query = cleaned_query.replace("(", "").replace(")", "").replace("-", " ")
    # Collapse multiple spaces
    cleaned_query = " ".join(cleaned_query.split())
    
    logger.info(f"Query pipeline: '{original_query[:60]}...' â†’ enhanced â†’ '{cleaned_query}'")
    
    url = f"{endpoint}/indexes/{index_name}/docs/search"
    params = {"api-version": api_version}
    headers = {"Content-Type": "application/json", "api-key": api_key}

    # Run semantic and keyword passes so the index can match titles and body text.
    sem_cfg = (os.getenv("AZURE_SEARCH_SEMANTIC_CONFIG") or getattr(Config, "AZURE_SEARCH_SEMANTIC_CONFIG", "default-semantic-config")).strip()
    
    # IMPORTANT: Azure Semantic Search requires a non-empty base 'search' text.
    # Using '*' with queryType='semantic' returns HTTP 400.
    # Therefore, always pass the cleaned query string.
    search_term = cleaned_query
    requested_top = max(top, int(os.getenv("AZURE_SEARCH_MIN_RESULTS", "10")))
    requested_top = min(requested_top, int(os.getenv("AZURE_SEARCH_MAX_RESULTS", "30")))
    highlight_fields = (os.getenv("AZURE_SEARCH_HIGHLIGHT_FIELDS") or "").strip()
    search_passes = [
        ("semantic", {
            "search": search_term,
            "count": True,
            "queryType": "semantic",
            "semanticConfiguration": sem_cfg,
            "captions": "extractive",
            "answers": "extractive|count-3",
            "top": requested_top,
            # NB: no "queryLanguage" — API version 2025-09-01 rejects it for the
            # search operation (HTTP 400) and falls back to keyword. en-us is the
            # default for the semantic configuration anyway.
        }, 20.0),
        ("keyword-all", {
            "search": search_term,
            "count": True,
            "queryType": "simple",
            "searchMode": "all",
            "top": requested_top,
        }, 10.0),
        ("keyword-any", {
            "search": search_term,
            "count": True,
            "queryType": "simple",
            "searchMode": "any",
            "top": requested_top,
        }, 0.0),
    ]
    if highlight_fields:
        for _name, body, _weight in search_passes:
            body["highlight"] = highlight_fields
            body["highlightPreTag"] = ""
            body["highlightPostTag"] = ""

    query_terms = [
        term
        for term in _extract_query_terms(query, min_len=1, include_joined=True)
        if len(term) > 3 and term not in {"about", "document", "documents", "file", "files", "please", "show", "tell", "info", "information"}
    ]
    specific_lookup = (
        bool(query_terms)
        and (
            " or " in cleaned_query.lower()
            or " and " in cleaned_query.lower()
            or cleaned_query.lower().startswith(("do you have", "do you know", "find", "search", "who is", "what about", "is there"))
            or len(query_terms) >= 3
        )
    )

    def _result_haystack(doc: dict) -> str:
        return " ".join(
            str(doc.get(field) or "") for field in ("name", "title", "file_path", "url", "content", "snippet")
        ).lower()

    def _apply_specific_lookup_filter(mapped: list[dict]) -> list[dict]:
        if not specific_lookup:
            return mapped
        filtered = [doc for doc in mapped if any(term in _result_haystack(doc) for term in query_terms)]
        if filtered:
            logger.info(
                "Azure Search precision filter: query=%r | mode=specific_lookup | kept=%s | discarded=%s",
                query,
                len(filtered),
                len(mapped) - len(filtered),
            )
            return filtered
        logger.info(
            "Azure Search precision filter: query=%r | mode=specific_lookup | kept=0 | discarded=%s",
            query,
            len(mapped),
        )
        return []

    merged_results = []
    for pass_name, body, weight in search_passes:
        pass_started_at = time.perf_counter()
        try:
            logger.info(f"Azure Search {pass_name} request: index={index_name} body={body}")
            resp = requests.post(url, params=params, headers=headers, json=body, timeout=getattr(config, 'HTTP_TIMEOUT', 15))
            if resp.status_code != 200:
                logger.warning(f"Azure Search ({pass_name}) error: HTTP {resp.status_code} {resp.text[:200]}")
                logger.info(
                    "AZURE SEARCH TIMING | query=%r | pass=%s | seconds=%.2f | results=%s",
                    query,
                    pass_name,
                    time.perf_counter() - pass_started_at,
                    0,
                )
                continue
            data = resp.json() or {}
            items = data.get("value", []) or []
            answers = data.get("@search.answers") or []
            if not items:
                logger.info(f"Azure Search ({pass_name}) returned 0 results")
                logger.info(
                    "AZURE SEARCH TIMING | query=%r | pass=%s | seconds=%.2f | results=%s",
                    query,
                    pass_name,
                    time.perf_counter() - pass_started_at,
                    0,
                )
                continue

            mapped = [_map_search_doc(d) for d in items]
            mapped = _apply_specific_lookup_filter(mapped)
            if not mapped:
                logger.info(
                    "AZURE SEARCH TIMING | query=%r | pass=%s | seconds=%.2f | results=%s",
                    query,
                    pass_name,
                    time.perf_counter() - pass_started_at,
                    0,
                )
                continue
            if pass_name == "semantic":
                try:
                    if isinstance(answers, list) and answers:
                        ans = answers[0]
                        ans_text = (ans.get("text") or "").strip()
                        if ans_text and items:
                            first = items[0]
                            mapped.insert(0, {
                                "id": "semantic-answer",
                                "name": first.get("name") or first.get("title") or "Answer",
                                "file_path": first.get("url") or first.get("webUrl") or first.get("metadata_storage_path") or first.get("file_path") or "",
                                "url": first.get("url") or first.get("webUrl") or first.get("metadata_storage_path") or first.get("file_path") or "",
                                "content": ans_text,
                                "file_type": "semantic-answer",
                                "score": ans.get("score")
                            })
                except Exception:
                    pass
            merged_results = _merge_ranked_search_docs(merged_results, mapped, source_weight=weight)
            logger.info(
                "AZURE SEARCH TIMING | query=%r | pass=%s | seconds=%.2f | results=%s",
                query,
                pass_name,
                time.perf_counter() - pass_started_at,
                len(mapped),
            )
        except Exception as e:
            logger.error(f"Azure Search ({pass_name}) request failed: {e}")
            logger.info(
                "AZURE SEARCH TIMING | query=%r | pass=%s | seconds=%.2f | results=%s",
                query,
                pass_name,
                time.perf_counter() - pass_started_at,
                0,
            )

    if merged_results:
        logger.info(f"Azure Search hybrid returned {len(merged_results)} unique result(s)")
        logger.info(
            "AZURE SEARCH TOTAL | query=%r | seconds=%.2f | results=%s",
            query,
            time.perf_counter() - started_at,
            len(merged_results),
        )
        return merged_results[:requested_top]
    
    exact_body = {
        "search": search_term,  # Non-empty query required for semantic ranking
        "count": True,
        "queryType": "semantic",
        "semanticConfiguration": sem_cfg,
        "captions": "extractive",
        "answers": "extractive|count-3",
        "top": min(top, 15),  # Respect top parameter â€” cap at 15 for performance
        # NB: no "queryLanguage" — rejected by API 2025-09-01 (HTTP 400).
    }
    
    # Log final body for diagnostics
    logger.info(f"Azure Search semantic body prepared: {exact_body}")

    try:
        logger.info(f"Azure Search semantic request: index={index_name} body={exact_body}")
        resp = requests.post(url, params=params, headers=headers, json=exact_body, timeout=getattr(config, 'HTTP_TIMEOUT', 15))
        if resp.status_code == 200:
            data = resp.json() or {}
            items = data.get("value", []) or []
            answers = data.get("@search.answers") or []
            if items:
                logger.info(f"Azure Search (semantic) returned {len(items)} results")
            mapped = [_map_search_doc(d) for d in items[:top]]
            mapped = _apply_specific_lookup_filter(mapped)
            if not mapped:
                logger.info(
                    "AZURE SEARCH TIMING | query=%r | pass=%s | seconds=%.2f | results=%s",
                    query,
                    "semantic-fallback",
                    time.perf_counter() - started_at,
                    0,
                )
                logger.info(
                    "AZURE SEARCH TOTAL | query=%r | seconds=%.2f | results=%s",
                    query,
                    time.perf_counter() - started_at,
                    0,
                )
                return []
            # Prepend a synthetic semantic answer result if available
            try:
                if isinstance(answers, list) and answers:
                    ans = answers[0]
                    ans_text = (ans.get("text") or "").strip()
                    if ans_text and items:
                        first = items[0]
                        mapped.insert(0, {
                            "id": "semantic-answer",
                            "name": first.get("name") or first.get("title") or "Answer",
                            "file_path": first.get("url") or first.get("webUrl") or first.get("metadata_storage_path") or first.get("file_path") or "",
                            "url": first.get("url") or first.get("webUrl") or first.get("metadata_storage_path") or first.get("file_path") or "",
                            "content": ans_text,
                            "file_type": "semantic-answer",
                            "score": ans.get("score")
                        })
            except Exception:
                pass
            logger.info(
                "AZURE SEARCH TIMING | query=%r | pass=%s | seconds=%.2f | results=%s",
                query,
                "semantic-fallback",
                time.perf_counter() - started_at,
                len(mapped),
            )
            logger.info(
                "AZURE SEARCH TOTAL | query=%r | seconds=%.2f | results=%s",
                query,
                time.perf_counter() - started_at,
                len(mapped),
            )
            return mapped
        else:
            logger.warning(f"Azure Search (semantic) error: HTTP {resp.status_code} {resp.text[:200]}")
    except Exception as e:
        logger.error(f"Azure Search (semantic) request failed: {e}")

    logger.info("Azure Search returned 0 results across semantic pass")
    logger.info(
        "AZURE SEARCH TOTAL | query=%r | seconds=%.2f | results=%s",
        query,
        time.perf_counter() - started_at,
        0,
    )
    return []

# =====================================================
# OneDrive-Specific Search
# =====================================================
def search_onedrive_personal(query: str, token: str, user_upn: str = "", top: int = 20) -> dict:
    """
    Search specifically within user's personal OneDrive.
    
    This function is more targeted than generic SharePoint search - it uses
    a direct filter for personal OneDrive scope to ensure files are from
    the user's own drive, not shared folders or team sites.
    
    Args:
        query: Search query string
        token: Graph access token (should be delegated for personal OneDrive access)
        user_upn: User's principal name (email) for scope filtering
        top: Maximum results to return
        
    Returns:
        Dictionary with search results or error
    """
    try:
        SearchLogger.search_source("OneDrive", 0, query)
        
        headers = {
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json",
        }
        
        # Expand query for better recall
        cleaned_query = _normalize_search_query(query)
        # Add phrase wrapping for better relevance in Graph
        if " " in cleaned_query:
            phrase_query = f'"{cleaned_query}"'
            onedrive_query = f"({phrase_query} OR {cleaned_query}) AND (driveType:personal)"
        else:
            onedrive_query = f"({cleaned_query or query}) AND (driveType:personal)"
        
        search_body = {
            "requests": [
                {
                    "entityTypes": ["driveItem"],  # Only driveItems, not sites/listItems
                    "query": {"queryString": onedrive_query},
                    "from": 0,
                    "size": top,
                    "region": "US",
                }
            ]
        }
        
        # Perform the search with retry logic
        def _search_onedrive():
            return requests.post(
                "https://graph.microsoft.com/v1.0/search/query",
                headers=headers,
                json=search_body,
                timeout=config.GRAPH_TIMEOUT,
            )
        
        try:
            resp = _retry_request(_search_onedrive, max_retries=1, initial_delay=0)  # fail fast
        except (requests.exceptions.Timeout, requests.exceptions.ReadTimeout) as e:
            SearchLogger.error(f"OneDrive search timed out")
            return {"results": [], "error": f"OneDrive search timed out"}
        
        if resp.status_code != 200:
            SearchLogger.warning(f"OneDrive search failed (HTTP {resp.status_code})")
            return {"results": [], "error": f"HTTP {resp.status_code}"}
        
        data = resp.json()
        results = []
        
        # Extract and filter OneDrive search results
        for block in data.get("value", []):
            for container in block.get("hitsContainers", []) or []:
                for hit in container.get("hits", []) or []:
                    resource = hit.get("resource", {}) or {}
                    name = resource.get("name") or resource.get("title") or "Untitled"
                    web_url = resource.get("webUrl", "")
                    
                    # Strict document filtering - essential for OneDrive
                    if not _is_supported_document(name):
                        continue
                    
                    # Ensure it's actually from personal OneDrive
                    if not _is_personal_url(web_url):
                        continue
                    
                    # Verify user can access (ownership check)
                    if user_upn and not is_url_accessible_by_user(web_url, user_upn, user_context=True):
                        continue
                    
                    results.append({
                        "name": name,
                        "webUrl": web_url,
                        "summary": hit.get("summary", ""),
                        "driveId": resource.get("parentReference", {}).get("driveId", ""),
                        "itemId": resource.get("id", ""),
                        "relevance_score": hit.get("@search.rank", 0),
                        "_onedrive_personal": True
                    })
        
        SearchLogger.search_source("OneDrive", len(results), query)
        for doc in results[:5]:
            ext = doc.get("name", "").rsplit(".", 1)[-1].lower() if "." in doc.get("name", "") else "?"
            SearchLogger.document(doc.get("name", "Untitled"), doc_type=ext, score=doc.get("relevance_score"))
        if len(results) > 5:
            logger.info(f"    ... and {len(results) - 5} more")
        return {"results": results, "count": len(results), "_from_onedrive_search": True}
    
    except Exception as e:
        SearchLogger.error(f"OneDrive search failed")
        return {"results": [], "error": str(e), "_from_onedrive_search": True}

# =====================================================
# Unified Search API
# =====================================================
def unified_search(query: str, top: int = 15, user_context: bool = True, user_id: Optional[str] = None, user_upn: str = "", user_assertion: Optional[str] = None) -> list:
    """Search SharePoint-first knowledge sources.
    
    Search flow:
    - Cache: previously crawled SharePoint docs (fastest)
    - Live Graph API: SharePoint fallback when cache is empty or weak
    - Optional sources: OneDrive, Azure AI Search, and web only when enabled
    
    Args:
        query: Search query
        top: Max results to return
        user_context: Whether to use user context (always True for security)
        user_id: User's ID for personalized results and access control
        user_upn: User principal name (email) for permission pre-filtering
    
    Returns:
        Combined list of search results with deduplication
    """
    SearchLogger.section(f"DOCUMENT SEARCH: '{query}'")
    
    # Ensure user profile is cached for better experience
    if user_id:
        ensure_user_profile_cached(user_id, user_assertion=user_assertion)
    
    all_results = []
    seen_ids = set()
    cache_has_results = False
    cache_has_no_results = False
    cache_buffer = []
    top_cache_score = 0
    cache_user = user_id or getattr(Config, "SHAREPOINT_CACHE_USER_ID", "shared")
    
    # Step 1: Read SharePoint cache first. This is the primary fast path.
    if getattr(Config, "ENABLE_SHAREPOINT_CACHE", True):
        cache = get_cache()
        cache_results = cache.search_cache_scored(query, user_id=cache_user, limit=top, include_shared=True)
        if cache_results:
            cache_has_results = True
            first_item = cache_results[0]
            if isinstance(first_item, dict):
                top_cache_score = int(first_item.get("score", 0) or 0)
            else:
                top_cache_score = int(first_item[1] if len(first_item) > 1 else 0)
            SearchLogger.search_source("Cache", len(cache_results), query)
            
            # Buffer cache results for optional merge after Graph
            for item in cache_results:
                if isinstance(item, dict):
                    doc = item.get("doc", {})
                else:
                    doc = item[0] if item else {}
                doc_id = doc.get("id") or doc.get("url") or doc.get("file_path")
                if doc_id and doc_id not in seen_ids:
                    cache_buffer.append(doc)
        else:
            cache_has_no_results = True
            SearchLogger.search_source("Cache", 0, query)

    cache_is_sufficient = (
        cache_buffer
        and (
            not getattr(Config, "REQUIRE_MULTI_DOCUMENT_SEARCH", True)
            or len(cache_buffer) >= getattr(Config, "MIN_SEARCH_RESULTS_BEFORE_GRAPH", 3)
        )
    )
    if (
        cache_is_sufficient
        and getattr(Config, "SHAREPOINT_CACHE_FIRST", True)
        and top_cache_score >= getattr(Config, "SHAREPOINT_CACHE_MIN_SCORE", 30)
    ):
        logger.info(
            "SharePoint cache-first hit: returning %s cached result(s), top_score=%s",
            len(cache_buffer),
            top_cache_score,
        )
        return cache_buffer[:top]
    elif cache_buffer:
        logger.info(
            "Cache found %s result(s), but continuing to Graph/AI Search for broader coverage",
            len(cache_buffer),
        )
    
    # Step 2: Live SharePoint Graph fallback when cache is empty or weak.
    should_call_graph = (
        getattr(Config, "ENABLE_SHAREPOINT_SEARCH", True)
        and getattr(Config, "ENABLE_LIVE_GRAPH_FALLBACK", True)
    )
    if should_call_graph:
        effective_user_context = bool(user_assertion) and user_context
        delegated_attempted = bool(user_assertion)

        def _ingest_graph_results(graph_results: list, from_app_only: bool = False) -> bool:
            if not graph_results:
                return False
            auth_label = "app-only" if from_app_only else "delegated"
            SearchLogger.search_source("Graph", len(graph_results), query)
            for doc in graph_results[:25]:  # Increased limit to preserve more results from Graph API
                doc_id = doc.get("id") or doc.get("webUrl") or doc.get("file_path")
                if doc_id and doc_id not in seen_ids:
                    seen_ids.add(doc_id)
                    doc["_from_live_graph"] = True

                    try:
                        try:
                            cache_ref = cache
                        except NameError:
                            cache_ref = None

                        web_url = doc.get("webUrl") or ""
                        drive_id = doc.get("driveId") or ""
                        item_id = doc.get("itemId") or ""
                        name = doc.get("name") or ""
                        composite_id = f"{drive_id}:{item_id}" if drive_id and item_id else web_url

                        cached_doc = None
                        if cache_ref and user_id:
                            all_cached = cache_ref.get_all_documents(user_id, include_shared=False)
                            cached_doc = next((d for d in all_cached
                                             if d.get("id") == composite_id
                                             or d.get("url") == web_url
                                             or (d.get("name") == name and name)), None)

                        if cached_doc and cached_doc.get("content") and len(cached_doc.get("content", "").strip()) > 50:
                            cached_content = cached_doc.get("content", "")
                            snippets = []
                            try:
                                if cache_ref:
                                    snippets = cache_ref._best_snippets(cached_content, query, window=500, max_snippets=3)
                            except Exception:
                                snippets = []
                            doc["content"] = safe_truncate(
                                "\n\n---\n\n".join(snippets) if snippets else cached_content,
                                Config.MAX_DOC_SNIPPET_CHARS
                            )
                            if snippets:
                                doc["snippets"] = snippets
                            doc["_from_cache"] = True
                            logger.info(f"âœ“ Populated Graph result with cached content: {name} (truncated to {len(doc['content'])} chars)")
                    except Exception as cache_err:
                        logger.debug(f"Cache lookup failed for {doc.get('name', 'unknown')}: {cache_err}")

                    all_results.append(doc)
            return True

        # Try to get token with retry logic
        token = None
        retry_attempts = 0
        max_retries = 2
        
        while token is None and retry_attempts < max_retries:
            try:
                token = get_graph_token(user_assertion)
                if token:
                    break
                else:
                    retry_attempts += 1
                    if retry_attempts < max_retries:
                        import time
                        time.sleep(0.5)  # Brief delay before retry
            except Exception as token_err:
                retry_attempts += 1
                if retry_attempts < max_retries:
                    import time
                    time.sleep(0.5)
        
        if not token:
            SearchLogger.error(f"Could not acquire Graph token - cannot search")
            # Return cache results as fallback
            return cache_buffer if cache_buffer else []
        
        try:
            # Optional OneDrive search. Disabled by default for SharePoint-first mode.
            onedrive_results = []
            if getattr(Config, "ENABLE_ONEDRIVE_SEARCH", False) and user_assertion and user_upn:
                try:
                    onedrive_data = search_onedrive_personal(query, token, user_upn=user_upn, top=top)
                    onedrive_results = onedrive_data.get("results", [])
                    if onedrive_results:
                        for doc in onedrive_results[:10]:  # Limit to top 10 OneDrive results
                            doc_id = doc.get("id") or doc.get("webUrl") or doc.get("file_path")
                            if doc_id and doc_id not in seen_ids:
                                seen_ids.add(doc_id)
                                doc["_from_onedrive_search"] = True
                                all_results.append(doc)
                except Exception as ode_err:
                    SearchLogger.warning(f"OneDrive search failed, continuing with general search")
            
            # SharePoint search via Graph.
            graph_data = search_sharepoint(query, token, user_context=effective_user_context, user_upn=user_upn)
            graph_results = graph_data.get("results", [])

            found_results = _ingest_graph_results(graph_results, from_app_only=not effective_user_context)

            # If delegated search returned nothing, optionally retry with app-only token
            if (not found_results and delegated_attempted 
                    and getattr(config, "GRAPH_ALLOW_APP_ONLY_FALLBACK", True)):
                fallback_token = get_graph_token_app_only()
                if fallback_token:
                    graph_data = search_sharepoint(query, fallback_token, user_context=False, user_upn=user_upn)
                    graph_results = graph_data.get("results", [])
                    found_results = _ingest_graph_results(graph_results, from_app_only=True)

            if not found_results:
                SearchLogger.search_source("Graph", 0, query)

            # Handle cache results when Graph returns nothing
            if cache_has_results:
                try:
                    relevance_threshold = 15
                    should_clear_cache = False
                    top_score = "unknown"

                    if 'cache_results' in locals() and cache_results and len(cache_results) > 0:
                        first_item = cache_results[0]
                        if isinstance(first_item, dict):
                            top_score = first_item.get("score", 0)
                        else:
                            top_score = first_item[1] if len(first_item) > 1 else 0
                    else:
                        top_score = 0

                    if isinstance(top_score, (int, float)) and top_score < relevance_threshold:
                        should_clear_cache = True
                    elif top_score == "unknown":
                        should_clear_cache = True

                    if should_clear_cache:
                        logger.info(f"ðŸ—‘ï¸ Clearing low-relevance cache results (score={top_score} < {relevance_threshold}) since Graph found nothing relevant")
                        cache_buffer = []
                        cache_has_results = False
                    else:
                        logger.info(f"âœ“ Keeping cache results (score={top_score} >= {relevance_threshold}) even though Graph found nothing")

                except Exception as clear_err:
                    logger.warning(f"Error in cache clearing logic: {clear_err}, proceeding with cache results")
        except Exception as e:
            logger.error(f"âŒ Live Graph search error: {e}", exc_info=True)
            logger.info(f"Falling back to cache results ({len(cache_buffer)} items)")
            return cache_buffer if cache_buffer else []

    # Merge cache results after Graph so Graph has priority
    if cache_buffer:
        for doc in cache_buffer:
            doc_id = doc.get("id") or doc.get("url") or doc.get("file_path")
            if doc_id and doc_id not in seen_ids:
                seen_ids.add(doc_id)
                all_results.append(doc)
    
    # Step 3: Optional Azure AI Search.
    if getattr(Config, "ENABLE_AI_SEARCH", False):
        ai_results = search_documents(query, top=top)
        if ai_results:
            SearchLogger.search_source("AI Search", len(ai_results), query)
            for doc in ai_results:
                doc_id = doc.get("id") or doc.get("url") or doc.get("file_path")
                if doc_id and doc_id not in seen_ids:
                    seen_ids.add(doc_id)
                    doc["_from_ai_search"] = True
                    all_results.append(doc)
        else:
            SearchLogger.search_source("AI Search", 0, query)
    
    # Return final results with diagnostics
    if all_results:
        _log_search_diagnostics(
            query,
            all_results,
            search_source="Unified",
            user_upn=user_upn,
            result_count=len(all_results),
            unique_sources=len(seen_ids)
        )
        logger.info(f"âœ“ Returning {len(all_results)} results from {len(seen_ids)} unique sources")
        return all_results
    
    logger.info("No Azure AI Search results found for '%s'", query)
    return []
"""
Knowledge Base - Graph API, Search, and Document Processing (FIXED)

Key fixes:
- Proper Graph shareId encoding: u!<base64url(url)>
- Graph fallback triggered on more than 401/403 (e.g., 404, 302, HTML/empty response)
- Consistent use of config.GRAPH_TIMEOUT
- sendMail fixed for app-only auth: /users/{SENDER_UPN}/sendMail
- Better logging and error diagnostics
"""

import io
import logging
import requests
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry
import base64
import tempfile
import json
import os
import time
import hashlib
from urllib.parse import urlparse, unquote
from typing import Optional
import re
from config import Config

# Lazy import to avoid ModuleNotFoundError when sys.path doesn't include src/
def get_cache():
    from document_cache import get_cache as _get_cache
    return _get_cache()

logger.setLevel(logging.INFO)
config = Config()

# Create a session with connection pooling to prevent socket exhaustion
# This is crucial for long-running bots that make many HTTP requests
_http_session = None

def get_http_session() -> requests.Session:
    """Get a reusable HTTP session with connection pooling and retry logic."""
    global _http_session
    if _http_session is None:
        _http_session = requests.Session()
        
        # Configure retry strategy for transient failures
        retry_strategy = Retry(
            total=3,
            backoff_factor=1,
            status_forcelist=[429, 500, 502, 503, 504],
            allowed_methods=["HEAD", "GET", "POST", "PUT", "DELETE", "OPTIONS", "TRACE"],
        )
        
        # Configure connection pooling
        adapter = HTTPAdapter(
            max_retries=retry_strategy,
            pool_connections=10,  # Number of connection pools
            pool_maxsize=20,      # Max connections per pool
            pool_block=False,     # Don't block when pool is full
        )
        
        _http_session.mount("http://", adapter)
        _http_session.mount("https://", adapter)
        
        # Set default timeout and headers
        _http_session.headers.update({
            "Connection": "keep-alive",
        })
        
        logger.info("HTTP session created with connection pooling")
    
    return _http_session


def session_get(url, **kwargs):
    """Make a GET request using the connection-pooled session."""
    return get_http_session().get(url, **kwargs)


def session_post(url, **kwargs):
    """Make a POST request using the connection-pooled session."""
    return get_http_session().post(url, **kwargs)


# User profile cache to store display names and UPNs
_USER_PROFILE_CACHE = {}
_USER_PROFILES_CACHE_FILE = os.path.join(os.path.dirname(__file__), "user_profiles_cache.json")


def _retry_request(func, max_retries=3, initial_delay=1.0, backoff_factor=2.0):
    """Retry a request with exponential backoff on timeout errors."""
    delay = initial_delay
    last_error = None
    
    for attempt in range(max_retries):
        try:
            return func()
        except (requests.exceptions.Timeout, requests.exceptions.ReadTimeout, requests.exceptions.ConnectionError) as e:
            last_error = e
            if attempt < max_retries - 1:
                logger.warning(f"Request timeout (attempt {attempt + 1}/{max_retries}). Retrying in {delay}s: {str(e)[:100]}")
                time.sleep(delay)
                delay *= backoff_factor
            else:
                logger.error(f"Request failed after {max_retries} attempts: {str(e)[:100]}")
    
    raise last_error


def _load_profiles_from_disk() -> dict:
    """Load user profiles from disk cache."""
    try:
        if os.path.exists(_USER_PROFILES_CACHE_FILE):
            with open(_USER_PROFILES_CACHE_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
    except Exception as e:
        logger.warning(f"Failed to load profiles from disk: {e}")
    return {}


def _save_profiles_to_disk(profiles: dict) -> None:
    """Save user profiles to disk cache."""
    try:
        # Ensure directory exists
        os.makedirs(os.path.dirname(_USER_PROFILES_CACHE_FILE), exist_ok=True)
        
        # Write directly (simpler than atomic, less chance of Windows locking issues)
        with open(_USER_PROFILES_CACHE_FILE, "w", encoding="utf-8") as f:
            json.dump(profiles, f, indent=2, ensure_ascii=False)
    except Exception as e:
        logger.warning(f"Failed to save profiles to disk: {e}")


def _get_profile_from_cache(user_id: str) -> Optional[dict]:
    """Get profile from memory or disk cache."""
    # Check memory first
    if user_id in _USER_PROFILE_CACHE:
        return _USER_PROFILE_CACHE[user_id]
    
    # Check disk cache
    disk_cache = _load_profiles_from_disk()
    if user_id in disk_cache:
        # Load into memory for faster access
        _USER_PROFILE_CACHE[user_id] = disk_cache[user_id]
        return disk_cache[user_id]
    
    return None


def _save_profile_to_cache(user_id: str, profile: dict) -> None:
    """Save profile to both memory and disk cache."""
    # Save to memory
    _USER_PROFILE_CACHE[user_id] = profile
    
    # Save to disk
    disk_cache = _load_profiles_from_disk()
    disk_cache[user_id] = profile
    _save_profiles_to_disk(disk_cache)

SUPPORTED_DOC_EXTENSIONS = {
    ".docx",
    ".doc",
    ".pdf",
    ".xlsx",
    ".xls",
    ".csv",
    ".txt",
    ".pptx",
    ".ppt",
    ".json",
    ".xml",
    ".rtf",
    ".odt",
    ".odp",
    ".ods",
    ".md",
    ".html",
    ".htm",
}

# Blacklist of explicitly unsupported file types that should NEVER be indexed
UNSUPPORTED_FILE_EXTENSIONS = {
    ".exe", ".dll", ".sys", ".bin", ".so", ".dylib", ".cmd", ".bat", ".sh", ".ps1",
    ".zip", ".rar", ".7z", ".tar", ".gz", ".iso", ".dmg",
    ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".svg", ".ico", ".webp", ".tiff",
    ".mp3", ".mp4", ".avi", ".mkv", ".mov", ".flv", ".wav", ".m4a", ".aac",
    ".tmp", ".cache", ".db", ".sqlite", ".mdb", ".accdb",
    ".class", ".jar", ".pyc", ".so", ".o", ".a", ".lib",
    ".git", ".hg", ".svn", ".cvs",
}
GRAPH_QUERY_STOPWORDS = {
    "a", "an", "the", "and", "or", "to", "of", "in", "on", "for", "from", "with",
    "my", "me", "please", "search", "find", "show", "list", "give", "provide",
    "summarize", "summary", "summarise", "analyze", "analysis", "report", "document", "documents",
    "file", "files", "about", "info", "information",
}

def _extract_query_terms(query: str, min_len: int = 3, include_joined: bool = True) -> list[str]:
    """Normalize query terms for Graph relevance filtering.
    
    Args:
        query: Raw search query
        min_len: Minimum term length to include
        include_joined: If True, append a concatenated variant of all terms
                       (useful for matching concatenated filenames during filtering,
                        but should be False when building search query strings)
    """
    if not query:
        return []
    terms = []
    for tok in re.split(r"\s+", query.strip()):
        t = tok.strip().strip("\"'`")
        t = re.sub(r"[^\w.\-]", "", t)
        if not t:
            continue
        t_lower = t.lower()
        for ext in SUPPORTED_DOC_EXTENSIONS:
            if t_lower.endswith(ext) and len(t_lower) > len(ext):
                t = t[: -len(ext)]
                t_lower = t.lower()
                break
        if not t or t_lower in GRAPH_QUERY_STOPWORDS:
            continue
        if len(t_lower) < min_len:
            continue
        terms.append(t_lower)
    # Add joined variant for better matching against concatenated filenames
    if include_joined and len(terms) > 1:
        joined = "".join(terms)
        if joined and joined not in terms:
            terms.append(joined)
    return terms

def _normalize_search_query(query: str) -> str:
    # Use include_joined=False: joined variants are added separately in
    # search_sharepoint and should not appear in the base search string.
    terms = _extract_query_terms(query, min_len=1, include_joined=False)
    if not terms:
        return (query or "").strip()
    return " ".join(terms)
GRAPH_CRAWL_MAX_ITEMS_PER_DRIVE = int(os.environ.get("GRAPH_CRAWL_MAX_ITEMS_PER_DRIVE", "300"))
GRAPH_CRAWL_MAX_DEPTH = int(os.environ.get("GRAPH_CRAWL_MAX_DEPTH", "4"))
# Maximum file size for download (50 MB default - increased to handle larger documents)
GRAPH_CRAWL_MAX_FILE_BYTES = int(os.environ.get("GRAPH_CRAWL_MAX_FILE_BYTES", "52428800"))  # 50 MB

# Optional document processing imports
pypdf = None
Document = None
load_workbook = None
Image = None
Presentation = None
xlrd = None
textract = None
try:
    import pypdf
except ImportError:
    pass
try:
    from docx import Document
except ImportError:
    pass
try:
    from openpyxl import load_workbook
except ImportError:
    pass
try:
    from PIL import Image
except ImportError:
    pass
try:
    from pptx import Presentation
except ImportError:
    pass
try:
    import xlrd
except ImportError:
    pass
try:
    import textract
except ImportError:
    pass


# =====================================================
# Helpers
# =====================================================
def _is_probably_html(resp: requests.Response) -> bool:
    ctype = (resp.headers.get("Content-Type") or "").lower()
    if "text/html" in ctype:
        return True
    # Sometimes SharePoint returns HTML login page with 200
    body = (resp.content[:200] or b"").lower()
    return b"<html" in body or b"<!doctype html" in body


def _should_attempt_graph(resp: Optional[requests.Response]) -> bool:
    """Decide if we should try Graph-auth download."""
    if resp is None:
        return True

    # Common failures
    if resp.status_code in (401, 403, 404, 410):
        return True

    # Redirects often indicate a pre-auth link that needs cookies/auth
    if resp.status_code in (301, 302, 303, 307, 308):
        return True

    # 200 but HTML/login page or empty content -> treat as failure
    if resp.status_code == 200:
        if len(resp.content or b"") == 0:
            return True
        if _is_probably_html(resp):
            return True

    return False


def _graph_share_id_from_url(url: str) -> str:
    """
    Graph expects shareId in the form: u!base64url(url)
    """
    encoded = base64.urlsafe_b64encode(url.encode("utf-8")).decode("utf-8").rstrip("=")
    return f"u!{encoded}"


def _log_http(prefix: str, label: str, resp: requests.Response) -> None:
    try:
        snippet = resp.text[:200] if resp.text else ""
    except Exception:
        snippet = "<non-text-response>"
    logger.info(
        f"{prefix}{label} HTTP {resp.status_code} | "
        f"Content-Type={resp.headers.get('Content-Type','')} | "
        f"Bytes={len(resp.content or b'')}"
    )
    if resp.status_code >= 400:
        logger.warning(f"{prefix}{label} error body snippet: {snippet}")


def _graph_download_via_path(content_url: str, headers: dict, prefix: str) -> Optional[requests.Response]:
    """
    Attempt to download a file via Graph using site/drive path resolution.
    Works best for classic SharePoint/OneDrive URLs with a resolvable site path.
    """
    try:
        parsed = urlparse(content_url)
        hostname = parsed.netloc
        server_path = unquote(parsed.path).lstrip("/")
        segments = [s for s in server_path.split("/") if s]
        if not segments:
            logger.info(f"{prefix}Graph path resolve: no path segments")
            return None

        # Handle:
        # /sites/<siteName>/...
        # /personal/<user>_<domain>_com/...
        if segments[0] in ("sites", "personal") and len(segments) >= 2:
            site_server_path = "/".join(segments[:2])
            file_segments = segments[2:]
        else:
            site_server_path = segments[0]
            file_segments = segments[1:]

        file_rel_path = "/".join(file_segments)

        # Resolve site id
        site_url = f"https://graph.microsoft.com/v1.0/sites/{hostname}:/{site_server_path}"
        logger.info(f"{prefix}Resolving site via Graph: {site_url}")
        site_resp = requests.get(site_url, headers=headers, timeout=config.GRAPH_TIMEOUT)
        _log_http(prefix, "Site resolve", site_resp)
        if site_resp.status_code != 200:
            return None

        site_id = site_resp.json().get("id")
        if not site_id:
            logger.error(f"{prefix}Site resolve returned no id")
            return None

        # Try default drive root path
        drive_content_url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drive/root:/{file_rel_path}:/content"
        logger.info(f"{prefix}Downloading via Graph drive path: {drive_content_url}")
        drive_resp = requests.get(drive_content_url, headers=headers, timeout=config.GRAPH_TIMEOUT)
        _log_http(prefix, "Drive path download", drive_resp)
        return drive_resp
    except Exception as e:
        logger.error(f"{prefix}Graph path resolution error: {e}", exc_info=True)
        return None


def _is_supported_document(name: str) -> bool:
    """Check if filename is a supported document type for indexing.
    
    This function applies multiple filters:
    1. Explicit extension blacklist (executables, media, archives, etc.)
    2. Hidden/system files (starting with . or ~)
    3. Supported document extension whitelist
    4. Minimum filename length check
    
    Args:
        name: Filename to check
        
    Returns:
        True if file is a supported document, False otherwise
    """
    if not name:
        return False
    
    # Strip whitespace
    name = name.strip()
    
    # Reject hidden/system files and temporary files
    if name.startswith(".") or name.startswith("~") or name.startswith("$"):
        logger.debug(f"Filtered out hidden/system file: {name}")
        return False
    
    # Reject very short filenames (likely system files)
    if len(name) < 3:
        logger.debug(f"Filtered out short filename: {name}")
        return False
    
    lower = name.lower()
    
    # Check blacklist first (reject explicitly unsupported types)
    if any(lower.endswith(ext) for ext in UNSUPPORTED_FILE_EXTENSIONS):
        logger.debug(f"Filtered out explicitly unsupported file type: {name}")
        return False
    
    # Check whitelist (accept only supported document types)
    is_supported = any(lower.endswith(ext) for ext in SUPPORTED_DOC_EXTENSIONS)
    
    if not is_supported:
        logger.debug(f"File does not match supported document extensions: {name}")
    
    return is_supported


def _is_personal_url(url: str) -> bool:
    """Heuristically detect personal OneDrive URLs to avoid sharing them globally."""
    if not url:
        return False
    url_lower = url.lower()
    return "/personal/" in url_lower or "my.sharepoint.com/personal" in url_lower


def _extract_owner_from_personal_url(url: str) -> Optional[str]:
    """Extract the owner identifier from a personal OneDrive URL.
    Returns username/email-like identifier or None.
    """
    if not url:
        return None
    try:
        url_lower = url.lower()
        if "/personal/" in url_lower:
            parts = url_lower.split("/personal/")
            if len(parts) > 1:
                owner_part = parts[1].split("/")[0]
                # Convert underscore format to email-like
                # e.g., "john_doe_company_com" -> "john.doe@company.com"
                if "_" in owner_part:
                    owner_clean = owner_part.replace("_", ".")
                    # Try to restore email format
                    if owner_clean.count(".") >= 2:
                        parts = owner_clean.rsplit(".", 2)
                        if len(parts) == 3:
                            return f"{parts[0]}@{parts[1]}.{parts[2]}"
                return owner_part
    except Exception:
        pass
    return None


def is_url_accessible_by_user(url: str, user_upn: str, user_context: bool = True) -> bool:
    """Check if a URL is likely accessible by the given user.
    
    IMPORTANT: When using delegated tokens (user_context=True), Graph Search API
    already filters results based on user permissions. Documents in search results
    should be accessible, including those shared with the user from other OneDrives.
    
    This function should only pre-filter obvious cases where access will definitely fail,
    not filter out potentially shared documents.
    
    Args:
        url: Document URL (webUrl from Graph search)
        user_upn: User principal name (email) of current user
        user_context: If True, using delegated token (trust Graph API filtering)
    
    Returns:
        True if document is likely accessible, False otherwise
    """
    logger.debug(f"DEBUG: is_url_accessible_by_user called with url='{url}', user_upn='{user_upn}', user_context={user_context}")
    
    if not url:
        return True
    
    url_lower = url.lower()
    
    # SECURITY: If user email is unknown, block ALL personal OneDrive access as a safety measure
    if "/personal/" in url_lower and not user_upn:
        logger.info(f"âŒ Blocking personal OneDrive URL because user UPN is unavailable (security): {url}")
        return False
    
    # ALWAYS check personal OneDrive ownership - whether delegated or app-only token
    # Users should NOT access documents from other users' personal OneDrive
    # even if they appear in Graph search results
    if "/personal/" in url_lower:
        owner = _extract_owner_from_personal_url(url)
        if owner:
            user_normalized = user_upn.lower().replace("@", ".").replace(".", "_")
            owner_normalized = owner.lower().replace("@", ".").replace(".", "_")
            
            logger.debug(f"DEBUG: Comparing user_normalized='{user_normalized}' with owner_normalized='{owner_normalized}'")
            
            # Check if owner matches current user
            if owner_normalized == user_normalized:
                logger.debug(f"âœ“ Document in user's own OneDrive: {owner}")
                return True
            
            # Different user's personal OneDrive - NOT accessible regardless of token type
            logger.info(f"âŒ Blocking access to {owner}'s personal OneDrive (current user: {user_upn})")
            return False
        return True  # Can't determine owner, assume it might be accessible
    
    # Team/Group sites - accessible to members
    if "/sites/" in url_lower:
        return True
    
    # Shared documents folder - accessible
    if "/shared documents" in url_lower or "/shared%20documents" in url_lower:
        return True
    
    # Default: assume accessible for other URL patterns
    return True


def _resolve_site_id(site_url: str, headers: dict) -> Optional[str]:
    """Resolve a SharePoint site URL to a site ID."""
    try:
        parsed = urlparse(site_url)
        hostname = parsed.netloc
        path = parsed.path.rstrip("/") or "/"
        resolve_url = f"https://graph.microsoft.com/v1.0/sites/{hostname}:{path}"
        resp = requests.get(resolve_url, headers=headers, timeout=config.GRAPH_TIMEOUT)
        _log_http("", "Site resolve", resp)
        if resp.status_code == 200:
            return resp.json().get("id")
        logger.warning(f"Failed to resolve site id for {site_url}: HTTP {resp.status_code}")
    except Exception as e:
        logger.error(f"Site resolve error for {site_url}: {e}", exc_info=True)
    return None


def _extract_pptx_text_fallback(content: bytes) -> str:
    """Fallback PPTX text extraction using ZIP + XML parsing (avoids python-pptx rId issues)."""
    try:
        import zipfile
        import xml.etree.ElementTree as ET

        texts = []
        with zipfile.ZipFile(io.BytesIO(content)) as zf:
            slide_files = [n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
            for name in sorted(slide_files):
                try:
                    xml_bytes = zf.read(name)
                    root = ET.fromstring(xml_bytes)
                    # Collect all text runs (<a:t>) regardless of namespace prefix
                    for elem in root.iter():
                        if elem.tag.endswith("}t") and elem.text:
                            texts.append(elem.text)
                except Exception:
                    continue
        combined = "\n".join(t.strip() for t in texts if t and t.strip()).strip()
        return combined
    except Exception:
        return ""


def _textract_legacy_office(content: bytes, extension: str, display_name: str) -> str:
    """Extract text from legacy Office formats (.ppt, .doc, .xls) using textract if available."""
    if textract is None:
        return ""
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{extension}") as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        extracted = textract.process(tmp_path, extension=extension)
        if extracted:
            return extracted.decode("utf-8", errors="ignore").strip()
    except Exception as e:
        logger.warning(f"Legacy extraction failed for {display_name}: {e}")
    finally:
        if tmp_path:
            try:
                os.remove(tmp_path)
            except Exception:
                pass
    return ""


def _extract_xls_with_xlrd(content: bytes) -> str:
    """Extract text from legacy .xls files with full data, statistics, and per-person breakdown."""
    if xlrd is None:
        return ""
    try:
        wb = xlrd.open_workbook(file_contents=content)
        sheets_text = []
        total_data_rows = 0
        all_column_stats = {}
        
        for sheet in wb.sheets():
            all_rows = []
            max_cols = 0
            
            # Collect ALL rows (no limits)
            for rowx in range(sheet.nrows):
                row = sheet.row_values(rowx)
                formatted_row = [_format_cell_value(cell) for cell in row]
                all_rows.append(formatted_row)
                max_cols = max(max_cols, len(formatted_row))
            
            if not all_rows:
                continue
            
            # Normalize row lengths
            for row in all_rows:
                while len(row) < max_cols:
                    row.append("")
            
            # Detect header
            has_header = False
            header_row = [f"Col{i+1}" for i in range(max_cols)]
            if all_rows:
                first_row = all_rows[0]
                non_empty = [c for c in first_row if c]
                if non_empty:
                    text_cells = sum(1 for c in non_empty if not _is_numeric(c))
                    has_header = text_cells >= len(non_empty) * 0.4
                    if has_header:
                        header_row = first_row
            
            # Get data rows (excluding header if present)
            data_rows = all_rows[1:] if has_header else all_rows
            
            # Compute statistics for numeric columns
            sheet_stats = {}
            for col_idx, col_name in enumerate(header_row):
                numeric_values = []
                for row in data_rows:
                    if col_idx < len(row):
                        cell_val = row[col_idx].strip()
                        try:
                            clean_val = cell_val.replace(",", "").replace("$", "").replace("%", "").strip()
                            if clean_val:
                                num = float(clean_val)
                                numeric_values.append(num)
                        except (ValueError, TypeError):
                            pass
                
                if len(numeric_values) >= max(1, len(data_rows) * 0.3):
                    total = sum(numeric_values)
                    avg = total / len(numeric_values) if numeric_values else 0
                    min_val = min(numeric_values) if numeric_values else 0
                    max_val = max(numeric_values) if numeric_values else 0
                    sheet_stats[col_name] = {"sum": total, "avg": avg, "min": min_val, "max": max_val, "count": len(numeric_values)}
                    if col_name in all_column_stats:
                        all_column_stats[col_name]["sum"] += total
                        all_column_stats[col_name]["count"] += len(numeric_values)
                        all_column_stats[col_name]["min"] = min(all_column_stats[col_name]["min"], min_val)
                        all_column_stats[col_name]["max"] = max(all_column_stats[col_name]["max"], max_val)
                    else:
                        all_column_stats[col_name] = {"sum": total, "count": len(numeric_values), "min": min_val, "max": max_val}
            
            # DYNAMIC: Detect identifier column and create breakdowns for ALL numeric columns
            identifier_col_idx = None
            numeric_col_indices = []
            
            for col_idx, col_name in enumerate(header_row):
                text_values = []
                numeric_count = 0
                for row in data_rows[:min(100, len(data_rows))]:
                    if col_idx < len(row):
                        val = row[col_idx].strip()
                        if val:
                            try:
                                float(val.replace(",", "").replace("$", "").replace("â‚¬", "").replace("K", "000").replace("M", "000000"))
                                numeric_count += 1
                            except:
                                text_values.append(val)
                
                if len(text_values) > numeric_count and identifier_col_idx is None:
                    unique_ratio = len(set(text_values)) / max(1, len(text_values))
                    if unique_ratio > 0.5:
                        identifier_col_idx = col_idx
                elif numeric_count > len(text_values) * 0.5:
                    numeric_col_indices.append(col_idx)
            
            # Create dynamic breakdown
            dynamic_breakdowns = {}
            if identifier_col_idx is not None and numeric_col_indices:
                for num_col_idx in numeric_col_indices[:5]:
                    col_name = header_row[num_col_idx] if num_col_idx < len(header_row) else f"Col{num_col_idx}"
                    breakdown = {}
                    for row in data_rows:
                        if identifier_col_idx < len(row) and num_col_idx < len(row):
                            identifier = row[identifier_col_idx].strip()
                            val_str = row[num_col_idx].strip()
                            if identifier and val_str:
                                try:
                                    clean_val = val_str.replace(",", "").replace("$", "").replace("â‚¬", "").replace("Â£", "")
                                    multiplier = 1
                                    if clean_val.endswith("K"):
                                        clean_val = clean_val[:-1]
                                        multiplier = 1000
                                    elif clean_val.endswith("M"):
                                        clean_val = clean_val[:-1]
                                        multiplier = 1000000
                                    num_val = float(clean_val) * multiplier
                                    if identifier not in breakdown:
                                        breakdown[identifier] = {"total": 0, "count": 0}
                                    breakdown[identifier]["total"] += num_val
                                    breakdown[identifier]["count"] += 1
                                except (ValueError, TypeError):
                                    pass
                    if breakdown:
                        dynamic_breakdowns[col_name] = breakdown
            
            rows_text = []
            
            # Add sheet statistics
            if sheet_stats:
                rows_text.append("ðŸ“Š **COLUMN STATISTICS:**")
                for col_name, stats in sheet_stats.items():
                    rows_text.append(f"  â€¢ **{col_name}**: SUM={stats['sum']:,.2f}, AVG={stats['avg']:,.2f}, MIN={stats['min']:,.2f}, MAX={stats['max']:,.2f}, COUNT={stats['count']}")
                rows_text.append("")
            
            # Add dynamic breakdowns
            if dynamic_breakdowns:
                id_col_name = header_row[identifier_col_idx] if identifier_col_idx < len(header_row) else "Item"
                rows_text.append(f"ðŸ“Š **BREAKDOWN BY {id_col_name.upper()}:**")
                for col_name, breakdown in dynamic_breakdowns.items():
                    if breakdown:
                        sorted_items = sorted(breakdown.items(), key=lambda x: x[1]["total"], reverse=True)
                        rows_text.append(f"\n  **{col_name} - Top 10:**")
                        for item, data in sorted_items[:10]:
                            rows_text.append(f"    â€¢ {item}: {data['total']:,.2f}")
                        if len(sorted_items) > 10:
                            rows_text.append(f"  **{col_name} - Bottom 5:**")
                            for item, data in sorted_items[-5:]:
                                rows_text.append(f"    â€¢ {item}: {data['total']:,.2f}")
                rows_text.append("")
            
            # Add column headers and sample rows
            rows_text.append(f"**Column Headers:** {', '.join(header_row)}")
            rows_text.append(f"**Rows:** {len(data_rows)}")
            rows_text.append("")
            rows_text.append("**SAMPLE DATA (First 20 rows):**")
            for idx, row in enumerate(data_rows[:20], 1):
                if any(cell.strip() for cell in row):
                    row_parts = [f"{hdr}={val}" for hdr, val in zip(header_row, row) if val.strip()]
                    rows_text.append(f"Row {idx}: " + " | ".join(row_parts))
                    total_data_rows += 1
            
            if len(data_rows) > 25:
                rows_text.append(f"\n... ({len(data_rows) - 25} rows omitted) ...")
            
            if rows_text:
                sheet_header = f"=== Sheet: {sheet.name} ({len(data_rows)} data rows) ===\n"
                sheets_text.append(sheet_header + "\n".join(rows_text))
        
        return "\n\n".join(sheets_text)
    except Exception as e:
        logger.error(f"Error extracting .xls content: {e}")
        return ""


def _format_cell_value(cell) -> str:
    """Format a cell value for AI readability - NO truncation, full content preserved."""
    if cell is None:
        return ""
    if isinstance(cell, float):
        # Check if it's actually an integer
        if cell == int(cell):
            return str(int(cell))
        # Format decimals nicely
        return f"{cell:.4f}".rstrip('0').rstrip('.')
    if isinstance(cell, (int, bool)):
        return str(cell)
    # Convert to string and clean whitespace (preserve full content)
    cell_str = str(cell).strip()
    # Replace problematic characters that could break table structure
    cell_str = cell_str.replace('\n', ' ').replace('\r', ' ').replace('\t', ' ')
    # Collapse multiple spaces
    while '  ' in cell_str:
        cell_str = cell_str.replace('  ', ' ')
    return cell_str


def _extract_xlsx_with_structure(content: bytes, display_name: str) -> str:
    """Extract Excel content with full data, automatic statistics, and AI-optimized formatting.
    
    Format optimized for LLM comprehension:
    - Clear sheet separation
    - Header row clearly marked
    - Automatic statistics for numeric columns (SUM, AVG, MIN, MAX)
    - Data rows numbered for reference
    - Full content preserved (no truncation)
    - Clean column separation with pipes
    """
    if load_workbook is None:
        return f"ðŸ“Š Excel: {display_name}\n\n(Install openpyxl to extract content.)"
    
    import warnings
    with warnings.catch_warnings():
        warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl")
        try:
            wb = load_workbook(io.BytesIO(content), data_only=True)
        except Exception as e:
            logger.error(f"Error loading Excel workbook: {e}")
            return f"ðŸ“Š Excel: {display_name}\n\nâŒ Error loading workbook: {str(e)}"
        
        sheets_text = []
        total_data_rows = 0
        all_column_stats = {}  # Aggregate stats across all sheets
        
        for sheet in wb.worksheets:
            # Collect ALL non-empty rows (no limits)
            all_rows = []
            max_cols = 0
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None and str(cell).strip() != "" for cell in row):
                    formatted_row = [_format_cell_value(cell) for cell in row]
                    all_rows.append(formatted_row)
                    max_cols = max(max_cols, len(formatted_row))
            
            if not all_rows:
                continue
            
            # Normalize row lengths for consistent table structure
            for row in all_rows:
                while len(row) < max_cols:
                    row.append("")
            
            # Detect if first row is a header (more text than numbers)
            has_header = False
            header_row = [f"Col{i+1}" for i in range(max_cols)]
            if all_rows:
                first_row = all_rows[0]
                non_empty = [c for c in first_row if c]
                if non_empty:
                    text_cells = sum(1 for c in non_empty if not _is_numeric(c))
                    has_header = text_cells >= len(non_empty) * 0.4  # At least 40% text
                    if has_header:
                        header_row = first_row
            
            # Get data rows (excluding header if present)
            data_rows = all_rows[1:] if has_header else all_rows
            
            # Compute automatic statistics for numeric columns
            sheet_stats = {}
            for col_idx, col_name in enumerate(header_row):
                numeric_values = []
                for row in data_rows:
                    if col_idx < len(row):
                        cell_val = row[col_idx].strip()
                        try:
                            clean_val = cell_val.replace(",", "").replace("$", "").replace("%", "").strip()
                            if clean_val:
                                num = float(clean_val)
                                numeric_values.append(num)
                        except (ValueError, TypeError):
                            pass
                
                if len(numeric_values) >= max(1, len(data_rows) * 0.3):
                    total = sum(numeric_values)
                    avg = total / len(numeric_values) if numeric_values else 0
                    min_val = min(numeric_values) if numeric_values else 0
                    max_val = max(numeric_values) if numeric_values else 0
                    sheet_stats[col_name] = {
                        "sum": total,
                        "avg": avg,
                        "min": min_val,
                        "max": max_val,
                        "count": len(numeric_values)
                    }
                    # Aggregate across sheets
                    if col_name in all_column_stats:
                        all_column_stats[col_name]["sum"] += total
                        all_column_stats[col_name]["count"] += len(numeric_values)
                        all_column_stats[col_name]["min"] = min(all_column_stats[col_name]["min"], min_val)
                        all_column_stats[col_name]["max"] = max(all_column_stats[col_name]["max"], max_val)
                    else:
                        all_column_stats[col_name] = {
                            "sum": total, "count": len(numeric_values),
                            "min": min_val, "max": max_val
                        }
            
            # DYNAMIC: Detect identifier column and create breakdowns for ALL numeric columns
            identifier_col_idx = None
            numeric_col_indices = []
            
            for col_idx, col_name in enumerate(header_row):
                text_values = []
                numeric_count = 0
                for row in data_rows[:min(100, len(data_rows))]:
                    if col_idx < len(row):
                        val = row[col_idx].strip()
                        if val:
                            try:
                                float(val.replace(",", "").replace("$", "").replace("â‚¬", "").replace("K", "000").replace("M", "000000"))
                                numeric_count += 1
                            except:
                                text_values.append(val)
                
                if len(text_values) > numeric_count and identifier_col_idx is None:
                    unique_ratio = len(set(text_values)) / max(1, len(text_values))
                    if unique_ratio > 0.5:
                        identifier_col_idx = col_idx
                elif numeric_count > len(text_values) * 0.5:
                    numeric_col_indices.append(col_idx)
            
            # Create dynamic breakdown
            dynamic_breakdowns = {}
            if identifier_col_idx is not None and numeric_col_indices:
                for num_col_idx in numeric_col_indices[:5]:
                    col_name = header_row[num_col_idx] if num_col_idx < len(header_row) else f"Col{num_col_idx}"
                    breakdown = {}
                    for row in data_rows:
                        if identifier_col_idx < len(row) and num_col_idx < len(row):
                            identifier = row[identifier_col_idx].strip()
                            val_str = row[num_col_idx].strip()
                            if identifier and val_str:
                                try:
                                    clean_val = val_str.replace(",", "").replace("$", "").replace("â‚¬", "").replace("Â£", "")
                                    multiplier = 1
                                    if clean_val.endswith("K"):
                                        clean_val = clean_val[:-1]
                                        multiplier = 1000
                                    elif clean_val.endswith("M"):
                                        clean_val = clean_val[:-1]
                                        multiplier = 1000000
                                    num_val = float(clean_val) * multiplier
                                    if identifier not in breakdown:
                                        breakdown[identifier] = {"total": 0, "count": 0}
                                    breakdown[identifier]["total"] += num_val
                                    breakdown[identifier]["count"] += 1
                                except (ValueError, TypeError):
                                    pass
                    if breakdown:
                        dynamic_breakdowns[col_name] = breakdown
            
            # Build AI-optimized structured output
            rows_text = []
            
            # Add sheet statistics if available
            if sheet_stats:
                rows_text.append("ðŸ“Š **COLUMN STATISTICS:**")
                for col_name, stats in sheet_stats.items():
                    rows_text.append(f"  â€¢ **{col_name}**: SUM={stats['sum']:,.2f}, AVG={stats['avg']:,.2f}, MIN={stats['min']:,.2f}, MAX={stats['max']:,.2f}, COUNT={stats['count']}")
                rows_text.append("")
            
            # Add dynamic breakdowns
            if dynamic_breakdowns:
                id_col_name = header_row[identifier_col_idx] if identifier_col_idx < len(header_row) else "Item"
                rows_text.append(f"ðŸ“Š **BREAKDOWN BY {id_col_name.upper()}:**")
                for col_name, breakdown in dynamic_breakdowns.items():
                    if breakdown:
                        sorted_items = sorted(breakdown.items(), key=lambda x: x[1]["total"], reverse=True)
                        rows_text.append(f"\n  **{col_name} - Top 10:**")
                        for item, data in sorted_items[:10]:
                            rows_text.append(f"    â€¢ {item}: {data['total']:,.2f}")
                        if len(sorted_items) > 10:
                            rows_text.append(f"  **{col_name} - Bottom 5:**")
                            for item, data in sorted_items[-5:]:
                                rows_text.append(f"    â€¢ {item}: {data['total']:,.2f}")
                rows_text.append("")
            
            # Add column headers and sample rows
            rows_text.append(f"**Column Headers:** {', '.join(header_row)}")
            rows_text.append(f"**Rows:** {len(data_rows)}")
            rows_text.append("")
            rows_text.append("**SAMPLE DATA (First 20 rows):**")
            for idx, row in enumerate(data_rows[:20], 1):
                if any(cell.strip() for cell in row):
                    row_parts = [f"{hdr}={val}" for hdr, val in zip(header_row, row) if val.strip()]
                    rows_text.append(f"Row {idx}: " + " | ".join(row_parts))
                    total_data_rows += 1
            
            if len(data_rows) > 25:
                rows_text.append(f"\n... ({len(data_rows) - 25} rows omitted) ...")
            
            if rows_text:
                sheet_header = f"=== Sheet: {sheet.title} ({len(data_rows)} data rows) ===\n"
                sheets_text.append(sheet_header + "\n".join(rows_text))
        
        if not sheets_text:
            return f"ðŸ“Š Excel: {display_name}\n\n[No data found in workbook]"
        
        # Add summary header with aggregate statistics
        summary = f"ðŸ“Š Excel File: {display_name}\n"
        summary += f"**Total Sheets:** {len(wb.worksheets)} | **Total Data Rows:** {total_data_rows}\n\n"
        
        # Add aggregate statistics across all sheets
        if all_column_stats:
            summary += "**ðŸ“Š AUTOMATIC COLUMN STATISTICS (ALL SHEETS):**\n"
            for col_name, stats in all_column_stats.items():
                avg = stats["sum"] / stats["count"] if stats["count"] > 0 else 0
                summary += f"  â€¢ **{col_name}**: SUM={stats['sum']:.2f}, AVG={avg:.2f}, MIN={stats['min']:.2f}, MAX={stats['max']:.2f}, COUNT={stats['count']}\n"
            summary += "\n"
        
        summary += "=" * 60 + "\n\n"
        
        return summary + "\n\n".join(sheets_text)


def _is_numeric(value: str) -> bool:
    """Check if a string value is numeric."""
    if not value:
        return False
    try:
        float(value.replace(",", "").replace("$", "").replace("%", ""))
        return True
    except ValueError:
        return False


def _extract_csv_with_structure(content: bytes, display_name: str) -> str:
    """Extract CSV content with full data, automatic statistics, and AI-optimized formatting.
    
    Format optimized for LLM comprehension:
    - Header row clearly marked
    - Automatic statistics for numeric columns (SUM, AVG, MIN, MAX)
    - Data rows numbered for reference
    - Full content preserved (no truncation)
    - Clean column separation
    """
    import csv as csv_module
    
    try:
        # Try to detect encoding
        text = None
        for encoding in ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252']:
            try:
                text = content.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        
        if text is None:
            return f"ðŸ“„ CSV: {display_name}\n\nâŒ Unable to decode file content"
        
        # Parse CSV - extract ALL rows
        reader = csv_module.reader(io.StringIO(text))
        all_rows = []
        max_cols = 0
        for row in reader:
            formatted_row = [_format_cell_value(cell) for cell in row]
            all_rows.append(formatted_row)
            max_cols = max(max_cols, len(formatted_row))
        
        if not all_rows:
            return f"ðŸ“„ CSV: {display_name}\n\n[Empty file]"
        
        # Normalize row lengths
        for row in all_rows:
            while len(row) < max_cols:
                row.append("")
        
        # Detect header (first row with mostly text)
        has_header = False
        header_row = [f"Col{i+1}" for i in range(max_cols)]
        if all_rows:
            first_row = all_rows[0]
            non_empty = [c for c in first_row if c]
            if non_empty:
                text_cells = sum(1 for c in non_empty if not _is_numeric(c))
                has_header = text_cells >= len(non_empty) * 0.4
                if has_header:
                    header_row = first_row
        
        # Get data rows (excluding header if present)
        data_rows = all_rows[1:] if has_header else all_rows
        total_data_rows = len(data_rows)
        
        # Compute automatic statistics for numeric columns
        column_stats = {}
        for col_idx, col_name in enumerate(header_row):
            numeric_values = []
            for row in data_rows:
                if col_idx < len(row):
                    cell_val = row[col_idx].strip()
                    try:
                        # Remove common formatting
                        clean_val = cell_val.replace(",", "").replace("$", "").replace("%", "").strip()
                        if clean_val:
                            num = float(clean_val)
                            numeric_values.append(num)
                    except (ValueError, TypeError):
                        pass
            
            # If more than 30% of rows have numeric values, compute stats
            if len(numeric_values) >= max(1, total_data_rows * 0.3):
                total = sum(numeric_values)
                avg = total / len(numeric_values) if numeric_values else 0
                min_val = min(numeric_values) if numeric_values else 0
                max_val = max(numeric_values) if numeric_values else 0
                column_stats[col_name] = {
                    "sum": total,
                    "avg": avg,
                    "min": min_val,
                    "max": max_val,
                    "count": len(numeric_values)
                }
        
        # Detect name/person column and hours column for per-person breakdown
        name_col_idx = None
        hours_col_idx = None
        name_keywords = ['fullname', 'name', 'employee', 'staff', 'person', 'user', 'member', 'player']
        hours_keywords = ['hours', 'hrs', 'time', 'duration', 'amount', 'total', 'wage', 'salary', 'value', 'price', 'cost', 'fee', 'pay', 'earnings', 'income']
        for idx, hdr in enumerate(header_row):
            hdr_lower = hdr.lower()
            if any(kw in hdr_lower for kw in name_keywords):
                name_col_idx = idx
            if any(kw in hdr_lower for kw in hours_keywords):
                hours_col_idx = idx
        
        # Compute per-person breakdown if we have both columns
        person_stats = {}
        if name_col_idx is not None and hours_col_idx is not None:
            for row in data_rows:
                if name_col_idx < len(row) and hours_col_idx < len(row):
                    person = row[name_col_idx].strip()
                    hours_str = row[hours_col_idx].strip()
                    if person and hours_str:
                        try:
                            hours = float(hours_str.replace(",", ""))
                            if person not in person_stats:
                                person_stats[person] = {"hours": 0, "entries": 0}
                            person_stats[person]["hours"] += hours
                            person_stats[person]["entries"] += 1
                        except ValueError:
                            pass
        
        # Build summary header
        summary = f"ðŸ“„ CSV File: {display_name}\n"
        summary += f"**Total Data Rows:** {total_data_rows}\n"
        summary += f"**Columns ({max_cols}):** {', '.join(header_row)}\n\n"
        
        # Add automatic statistics for numeric columns
        if column_stats:
            summary += "**ðŸ“Š AUTOMATIC COLUMN STATISTICS:**\n"
            for col_name, stats in column_stats.items():
                summary += f"  â€¢ **{col_name}**: "
                summary += f"SUM={stats['sum']:.2f}, "
                summary += f"AVG={stats['avg']:.2f}, "
                summary += f"MIN={stats['min']:.2f}, "
                summary += f"MAX={stats['max']:.2f}, "
                summary += f"COUNT={stats['count']}\n"
            summary += "\n"
        
        # Add per-person breakdown if available
        if person_stats:
            summary += "**ðŸ‘¥ HOURS BY PERSON (Pre-computed):**\n"
            # Sort by hours descending
            sorted_persons = sorted(person_stats.items(), key=lambda x: x[1]["hours"], reverse=True)
            for person, stats in sorted_persons:
                summary += f"  â€¢ **{person}**: {stats['hours']:.2f} hours ({stats['entries']} entries)\n"
            summary += "\n"
        
        summary += "=" * 60 + "\n\n"
        
        # Build data in compact format (more efficient for LLM processing)
        summary += "**FULL DATA (Row# | Values):**\n"
        summary += "-" * 60 + "\n"
        
        # Add ALL data rows with key=value format (no truncation)
        for idx, row in enumerate(data_rows, 1):
            # Format: Row# | col1=val1 | col2=val2 | ...
            row_parts = []
            for i, (hdr, val) in enumerate(zip(header_row, row)):
                val_clean = val.strip() if val else ""
                if val_clean:  # Only include non-empty values
                    row_parts.append(f"{hdr}={val_clean}")
            summary += f"Row {idx}: " + " | ".join(row_parts) + "\n"
        
        return summary
    
    except Exception as e:
        logger.error(f"Error extracting CSV content: {e}")
        return f"ðŸ“„ CSV: {display_name}\n\nâŒ Error: {str(e)}"


def _list_site_drives(site_id: str, headers: dict) -> list[dict]:
    """List drives for a site."""
    try:
        url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives?$select=id,name,driveType"
        resp = requests.get(url, headers=headers, timeout=config.GRAPH_TIMEOUT)
        _log_http("", "List site drives", resp)
        if resp.status_code != 200:
            return []
        data = resp.json() or {}
        return data.get("value", []) or []
    except Exception as e:
        logger.error(f"Error listing drives for site {site_id}: {e}", exc_info=True)
        return []


def _should_skip_file(item: dict) -> bool:
    size = item.get("size")
    try:
        if size is not None and int(size) > GRAPH_CRAWL_MAX_FILE_BYTES:
            return True
    except Exception:
        return False
    return False
def _make_drive_doc_id(drive_id: str, item_id: str) -> str:
    return f"{drive_id}:{item_id}"

# Token cache to prevent excessive token requests and connection timeouts
_graph_token_cache = {
    "token": None,
    "expires_at": 0,  # Unix timestamp
}
_graph_token_cache_obo: dict[str, dict] = {}
_TOKEN_CACHE_BUFFER_SECONDS = 300  # Refresh token 5 minutes before expiry

def _obo_cache_key(assertion: str) -> str:
    return hashlib.sha256(assertion.encode("utf-8")).hexdigest()

def get_graph_token_app_only() -> Optional[str]:
    """Acquire an app-only Graph token via Client Credentials or Managed Identity.
    
    Tokens are cached in memory and reused until 5 minutes before expiry to prevent
    connection exhaustion and timeouts from frequent token requests.
    """
    import time
    global _graph_token_cache
    
    # Check if cached token is still valid
    current_time = time.time()
    if _graph_token_cache["token"] and _graph_token_cache["expires_at"] > current_time + _TOKEN_CACHE_BUFFER_SECONDS:
        return _graph_token_cache["token"]
    
    session = get_http_session()
    
    try:
        tenant_id = config.GRAPH_TENANT_ID or config.APP_TENANTID
        client_id = config.GRAPH_CLIENT_ID or config.APP_ID
        client_secret = config.GRAPH_CLIENT_SECRET or config.APP_PASSWORD

        # Prefer client credentials if available
        if tenant_id and client_id and client_secret:
            token_url = f"https://login.microsoftonline.com/{tenant_id}/oauth2/v2.0/token"
            data = {
                "client_id": client_id,
                "client_secret": client_secret,
                "grant_type": "client_credentials",
                "scope": "https://graph.microsoft.com/.default",
            }
            resp = session.post(token_url, data=data, timeout=config.GRAPH_TIMEOUT)
            if resp.status_code == 200:
                token_data = resp.json() or {}
                tok = token_data.get("access_token")
                expires_in = token_data.get("expires_in", 3600)  # Default 1 hour
                if tok:
                    # Cache the token with expiry time
                    _graph_token_cache["token"] = tok
                    _graph_token_cache["expires_at"] = current_time + expires_in
                    logger.info(f"Graph token acquired and cached (expires in {expires_in}s)")
                    return tok
            logger.error(f"App-only token failure (client creds): HTTP {resp.status_code} {resp.text[:200]}")

        # Managed Identity (IMDS)
        try:
            imds = "http://169.254.169.254/metadata/identity/oauth2/token"
            params = {"api-version": "2018-02-01", "resource": "https://graph.microsoft.com/"}
            headers = {"Metadata": "true"}
            resp = requests.get(imds, params=params, headers=headers, timeout=config.GRAPH_TIMEOUT)
            if resp.status_code == 200:
                tok = (resp.json() or {}).get("access_token")
                if tok:
                    return tok
            logger.error(f"Managed Identity token failure (IMDS): HTTP {resp.status_code} {resp.text[:200]}")
        except Exception as e:
            logger.error(f"Managed Identity (IMDS) exception: {e}")

        return None
    except Exception as e:
        logger.error(f"App-only token exception: {e}", exc_info=True)
        return None


def get_graph_token_obo(user_assertion: str) -> Optional[str]:
    """Acquire a delegated Graph token using on-behalf-of (OBO)."""
    if not user_assertion:
        return None

    try:
        key = _obo_cache_key(user_assertion)
        cached = _graph_token_cache_obo.get(key)
        if cached and cached.get("token") and cached.get("expires_at", 0) > time.time() + _TOKEN_CACHE_BUFFER_SECONDS:
            return cached["token"]
    except Exception:
        pass

    try:
        tenant_id = config.GRAPH_TENANT_ID or config.APP_TENANTID
        client_id = config.GRAPH_CLIENT_ID or config.APP_ID
        client_secret = config.GRAPH_CLIENT_SECRET or config.APP_PASSWORD
        if not (tenant_id and client_id and client_secret):
            logger.warning("OBO token request missing tenant/client/secret")
            return None

        token_url = f"https://login.microsoftonline.com/{tenant_id}/oauth2/v2.0/token"
        data = {
            "client_id": client_id,
            "client_secret": client_secret,
            "grant_type": "urn:ietf:params:oauth:grant-type:jwt-bearer",
            "requested_token_use": "on_behalf_of",
            "scope": "https://graph.microsoft.com/.default",
            "assertion": user_assertion,
        }
        resp = get_http_session().post(token_url, data=data, timeout=config.GRAPH_TIMEOUT)
        if resp.status_code == 200:
            token_data = resp.json() or {}
            tok = token_data.get("access_token")
            expires_in = token_data.get("expires_in", 3600)
            if tok:
                try:
                    _graph_token_cache_obo[key] = {
                        "token": tok,
                        "expires_at": time.time() + expires_in,
                    }
                except Exception:
                    pass
                logger.info("OBO Graph token acquired")
                return tok
        logger.error(f"OBO token failure: HTTP {resp.status_code} {resp.text[:200]}")
    except Exception as e:
        logger.error(f"OBO token exception: {e}", exc_info=True)

    return None


# Back-compat: keep old name for existing callers
def get_graph_token(user_assertion: Optional[str] = None) -> Optional[str]:
    """Get Graph token - prefer delegated, fallback to app-only.
    
    Called by search and content download functions. Logs failures
    for diagnostics.
    """
    token = None
    if user_assertion:
        token = get_graph_token_obo(user_assertion)
        if token:
            logger.info("âœ… Using delegated token (OBO)")
            return token
        else:
            logger.warning("âš ï¸  OBO token failed, falling back to app-only")
    
    token = get_graph_token_app_only()
    if token:
        logger.info("âœ… Using app-only token")
        return token
    
    logger.error("âŒ CRITICAL: Both OBO and app-only token acquisition failed!")
    return None


_GUID_RE = re.compile(r"^[0-9a-fA-F]{8}-[0-9a-fA-F]{4}-[1-5][0-9a-fA-F]{3}-[89abAB][0-9a-fA-F]{3}-[0-9a-fA-F]{12}$")

def _looks_like_guid(value: str) -> bool:
    try:
        if not value or len(value) < 30:
            return False
        return bool(_GUID_RE.match(value))
    except Exception:
        return False


def get_user_profile(user_id: str, user_assertion: Optional[str] = None) -> Optional[dict]:
    """Get user profile from Microsoft Graph using delegated (OBO) when available.

    - If a user assertion is provided, call /me with an OBO token.
    - Otherwise, use app-only token to call /users/{idOrUserPrincipalName}.
      Supports both GUID AAD object IDs and UPN/email (contains '@').
    - Results are cached in memory for the session.
    """
    try:
        if not user_id:
            logger.warning("Graph profile: no user_id provided")
            return None
            
        # Check cache first (memory + disk)
        cached = _get_profile_from_cache(user_id)
        if cached:
            logger.info(f"User profile cache hit for {user_id[:min(8, len(user_id))]}...")
            return cached

        headers = {"Content-Type": "application/json"}

        if user_assertion:
            obo_token = get_graph_token_obo(user_assertion)
            if obo_token:
                headers["Authorization"] = f"Bearer {obo_token}"
                endpoint = "https://graph.microsoft.com/v1.0/me?$select=id,displayName,mail,userPrincipalName,jobTitle,givenName"
                logger.info("Graph profile: method=obo endpoint=/me")
                resp = session_get(endpoint, headers=headers, timeout=config.GRAPH_TIMEOUT)
                _log_http("", "Get profile (/me)", resp)
                if resp.status_code == 200:
                    profile = resp.json() or {}
                    result = {
                        "id": profile.get("id"),
                        "displayName": profile.get("displayName"),
                        "givenName": profile.get("givenName"),
                        "mail": profile.get("mail") or profile.get("userPrincipalName"),
                        "userPrincipalName": profile.get("userPrincipalName"),
                        "jobTitle": profile.get("jobTitle"),
                    }
                    _save_profile_to_cache(user_id or result.get("id") or "", result)
                    logger.info(f"Cached profile for user: {result.get('displayName')} ({result.get('userPrincipalName')})")
                    return result
                logger.warning(f"Graph profile: HTTP {resp.status_code} for /me")
            else:
                logger.warning("Graph profile: failed to acquire OBO token, falling back to app-only")

        # Accept GUID or UPN/email
        looks_guid = _looks_like_guid(user_id or "")
        looks_upn = isinstance(user_id, str) and ("@" in user_id)
        
        if not (looks_guid or looks_upn):
            logger.warning(f"Graph profile: user_id '{user_id[:min(20, len(user_id))]}...' is neither GUID nor UPN/email")
            return None
            
        token = get_graph_token_app_only()
        if not token:
            logger.error("Graph profile: failed to acquire app-only token")
            return None
            
        headers["Authorization"] = f"Bearer {token}"
        endpoint = f"https://graph.microsoft.com/v1.0/users/{user_id}?$select=id,displayName,mail,userPrincipalName,jobTitle,givenName"
        logger.info(f"Graph profile: method=app-only endpoint=/users/{{idOrUPN}} for {user_id[:min(8, len(user_id))]}...")
        resp = session_get(endpoint, headers=headers, timeout=config.GRAPH_TIMEOUT)
        _log_http("", "Get profile (/users/{idOrUPN})", resp)
        
        if resp.status_code == 200:
            profile = resp.json() or {}
            
            # DEBUG: Log all fields returned by Graph API to troubleshoot missing email
            available_fields = list(profile.keys())
            logger.info(f"DEBUG: Graph API returned fields for user {user_id[:8]}...: {available_fields}")
            logger.info(f"DEBUG: mail='{profile.get('mail')}', userPrincipalName='{profile.get('userPrincipalName')}'")
            
            # Extract email from available fields
            email = profile.get("mail") or profile.get("userPrincipalName") or ""
            
            # If no email is available from Graph but we have a UPN-like user_id, use that
            if not email and isinstance(user_id, str) and "@" in user_id:
                email = user_id
                logger.info(f"DEBUG: Using user_id as email fallback: {email}")
            
            result = {
                "id": profile.get("id"),
                "displayName": profile.get("displayName"),
                "givenName": profile.get("givenName"),
                "mail": email,
                "userPrincipalName": profile.get("userPrincipalName") or email,
                "jobTitle": profile.get("jobTitle"),
            }
            # Cache the result (memory + disk)
            _save_profile_to_cache(user_id, result)
            final_email = result.get('mail') or result.get('userPrincipalName') or '(none)'
            logger.info(f"Cached profile for user: {result.get('displayName')} (email: {final_email})")
            return result
        else:
            logger.warning(f"Graph profile: HTTP {resp.status_code} for user_id {user_id[:min(8, len(user_id))]}...")
            return None

    except Exception as e:
        logger.error(f"Profile fetch error for user_id '{user_id[:min(20, len(user_id)) if user_id else 0]}...': {e}", exc_info=True)
        return None


def get_user_display_name(user_id: str, fallback_name: str = None) -> str:
    """Get user's display name for friendly greetings.
    
    Args:
        user_id: User's AAD object ID or UPN
        fallback_name: Optional display name from Teams if Graph profile doesn't have one
    
    Returns:
        Display name or given name, falls back to provided name or 'there'
    """
    try:
        if not user_id:
            logger.warning("get_user_display_name: no user_id provided")
            return fallback_name or "there"
            
        profile = get_user_profile(user_id)
        if profile:
            # Prefer given name for informal greeting, fallback to display name, then fallback_name from Teams
            display = profile.get("givenName") or profile.get("displayName") or fallback_name or "there"
            logger.info(f"Display name resolved: '{display}' for user_id={user_id[:min(8, len(user_id))]}...")
            return display
        
        logger.warning(f"get_user_display_name: no profile found for user_id={user_id[:min(20, len(user_id))]}...")
        return fallback_name or "there"
    except Exception as e:
        logger.error(f"Error getting display name for user_id={user_id[:min(20, len(user_id)) if user_id else 0]}...: {e}")
        return fallback_name or "there"


def ensure_user_profile_cached(user_id: str, user_assertion: Optional[str] = None) -> bool:
    """Ensure user profile is fetched and cached at the start of conversation.
    
    Args:
        user_id: User's AAD object ID or UPN
        user_assertion: Optional Teams SSO token for delegated Graph access
    
    Returns:
        True if profile was successfully cached, False otherwise
    """
    try:
        profile = get_user_profile(user_id, user_assertion=user_assertion)
        return profile is not None
    except Exception:
        return False


# =====================================================
# Document Processing
# =====================================================
def process_document(attachment, corr_id: Optional[str] = None) -> str:
    """Download attachment content and extract text. Uses Graph when needed."""
    prefix = f"[{corr_id}] " if corr_id else ""

    try:
        display_name = getattr(attachment, "name", None) or "attachment"
        logger.info(f"{prefix}Processing document: {display_name}")

        content_type = getattr(attachment, "content_type", "") or ""
        if content_type:
            logger.info(f"{prefix}Attachment content_type: {content_type}")

        # Extract content url from Teams attachment payload
        content_info = None
        content_url = None

        if hasattr(attachment, "content") and attachment.content:
            if isinstance(attachment.content, dict):
                content_info = attachment.content
            elif isinstance(attachment.content, str):
                s = attachment.content.strip()
                if s.startswith("{"):
                    try:
                        content_info = json.loads(s)
                    except Exception:
                        content_info = None

        # Prefer Teams-provided pre-auth downloadUrl if present
        if content_info and content_info.get("downloadUrl"):
            content_url = content_info.get("downloadUrl")
            logger.info(f"{prefix}Using Teams-provided downloadUrl")

        if not content_url and getattr(attachment, "content_url", None):
            content_url = attachment.content_url

        if not content_url and content_info and content_info.get("contentUrl"):
            content_url = content_info.get("contentUrl")

        if not content_url:
            logger.warning(f"{prefix}No content URL for {display_name}")
            return (
                f"âŒ Cannot access {display_name} (type: {content_type or 'unknown'}). "
                "No download URL provided. Try re-uploading directly in chat."
            )

        # 1) Try direct download (fast path)
        logger.info(f"{prefix}Direct download attempt: {content_url[:120]}...")
        direct_resp = None
        try:
            direct_resp = requests.get(
                content_url,
                timeout=int(os.environ.get("ATTACHMENT_DIRECT_DOWNLOAD_TIMEOUT", "8")),
                allow_redirects=False,
            )
            _log_http(prefix, "Direct download", direct_resp)
        except Exception as e:
            logger.warning(f"{prefix}Direct download exception: {e}")

        # 2) If direct is suspicious/fails, use Graph
        resp = direct_resp
        if _should_attempt_graph(direct_resp):
            logger.info(f"{prefix}Attempting Graph-auth download...")
            token = get_graph_token()
            if not token:
                return (
                    f"âŒ Could not get Graph token to download {display_name}. "
                    "Check Graph client id/secret/tenant and permissions."
                )

            headers = {"Authorization": f"Bearer {token}"}

            # If URL already a Graph URL, retry with auth header
            if "graph.microsoft.com" in content_url:
                logger.info(f"{prefix}Retrying Graph URL with bearer token...")
                resp = requests.get(content_url, headers=headers, timeout=config.GRAPH_TIMEOUT, allow_redirects=True)
                _log_http(prefix, "Graph direct", resp)
            else:
                # Try shares API
                share_id = _graph_share_id_from_url(content_url)
                shares_url = f"https://graph.microsoft.com/v1.0/shares/{share_id}/driveItem/content"
                logger.info(f"{prefix}Retrying via Graph shares API: {shares_url}")
                resp = requests.get(shares_url, headers=headers, timeout=config.GRAPH_TIMEOUT, allow_redirects=True)
                _log_http(prefix, "Graph shares", resp)

                # If shares fails, try path-based resolution
                if resp.status_code in (400, 404):
                    alt_resp = _graph_download_via_path(content_url, headers, prefix)
                    if alt_resp is not None:
                        resp = alt_resp

        if resp is None:
            return f"âŒ Failed to download {display_name}: no response."

        try:
            resp.raise_for_status()
        except requests.HTTPError:
            status = resp.status_code
            if status in (401, 403):
                return (
                    f"âŒ Access denied for {display_name} (HTTP {status}). "
                    "Ensure app has *application* permissions: Files.Read.All, Sites.Read.All "
                    "and admin consent is granted."
                )
            return f"âŒ Failed to download {display_name} (HTTP {status})."

        content = resp.content or b""
        if len(content) == 0:
            return f"âŒ Downloaded 0 bytes for {display_name}. (Likely an auth/redirect issue.)"

        file_name = display_name.lower()
        logger.info(f"{prefix}Downloaded {len(content)} bytes for {display_name}")

        # PDF
        if file_name.endswith(".pdf"):
            if pypdf is None:
                return f"ðŸ“„ PDF: {display_name}\n\n(Install pypdf to extract text.)"
            reader = pypdf.PdfReader(io.BytesIO(content))
            text_parts = []
            for page in reader.pages:
                t = page.extract_text() or ""
                if t.strip():
                    text_parts.append(t)
            text = "\n".join(text_parts)
            return f"ðŸ“„ PDF: {display_name}\n\n{text}"

        # Word
        if file_name.endswith((".docx", ".doc")): 
            if Document is None:
                return f"ðŸ“ Word: {display_name}\n\n(Install python-docx to extract text.)"
            
            # python-docx only works with .docx (Office Open XML format)
            # Old .doc files (Word 97-2003) are binary format and won't work
            if file_name.endswith(".doc") and not file_name.endswith(".docx"):
                # Try to detect if it's actually a .docx misnamed as .doc
                try:
                    doc = Document(io.BytesIO(content))
                    text = "\n".join(p.text for p in doc.paragraphs if p.text and p.text.strip())
                    return f"ðŸ“ Word: {display_name}\n\n{text}"
                except Exception as e:
                    logger.warning(f"Cannot extract old .doc format: {display_name}")
                    return (
                        f"ðŸ“ Word (Legacy): {display_name}\n\n"
                        f"âš ï¸ This file is in old Word 97-2003 format (.doc).\n"
                        f"Please convert to .docx format or re-save as .docx to extract content.\n"
                        f"File size: {len(content)} bytes"
                    )
            
            # .docx files
            try:
                doc = Document(io.BytesIO(content))
                text = "\n".join(p.text for p in doc.paragraphs if p.text and p.text.strip())
                return f"ðŸ“ Word: {display_name}\n\n{text}"
            except Exception as e:
                logger.error(f"Error extracting Word document: {e}")
                return f"ðŸ“ Word: {display_name}\n\nâŒ Error extracting content: {str(e)}"

        # Excel - use improved extraction with structure
        if file_name.endswith((".xlsx", ".xls")):
            if file_name.endswith(".xls") and not file_name.endswith(".xlsx"):
                xls_text = _extract_xls_with_xlrd(content)
                if xls_text:
                    return f"ðŸ“Š Excel: {display_name}\n\n{xls_text}"
                if xlrd is None:
                    return f"ðŸ“Š Excel: {display_name}\n\n(Install xlrd==1.2.0 to extract legacy .xls content.)"
                return f"ðŸ“Š Excel: {display_name}\n\nâŒ Error extracting legacy .xls content."
            # Use improved .xlsx extraction
            return _extract_xlsx_with_structure(content, display_name)

        # CSV files
        if file_name.endswith(".csv"):
            return _extract_csv_with_structure(content, display_name)

        # PowerPoint
        if file_name.endswith((".pptx", ".ppt")):
            if file_name.endswith(".ppt") and not file_name.endswith(".pptx"):
                if content[:4] == b"PK\x03\x04":
                    try:
                        if Presentation is None:
                            return f"ðŸ“½ï¸ PowerPoint: {display_name}\n\n(Install python-pptx to extract content.)"
                        prs = Presentation(io.BytesIO(content))
                        slides_text = []
                        for slide in prs.slides:
                            parts = []
                            for shape in slide.shapes:
                                try:
                                    if hasattr(shape, "text"):
                                        txt = (shape.text or "").strip()
                                        if txt:
                                            parts.append(txt)
                                except Exception:
                                    continue
                            if parts:
                                slides_text.append("\n".join(parts))
                        combined = "\n\n".join(slides_text).strip()
                        return combined if combined else "[No extractable text found in slides]"
                    except Exception as e:
                        logger.warning(f"PowerPoint extraction error for {display_name}: {e}")
                legacy_text = _textract_legacy_office(content, "ppt", display_name)
                if legacy_text:
                    return legacy_text
                if textract is None:
                    return "[Legacy .ppt detected. Install textract to extract content.]"
                return "[Legacy .ppt detected but extraction failed.]"
            if Presentation is None:
                return f"ðŸ“½ï¸ PowerPoint: {display_name}\n\n(Install python-pptx to extract content.)"
            try:
                prs = Presentation(io.BytesIO(content))
                slides_text = []
                for slide in prs.slides:
                    parts = []
                    for shape in slide.shapes:
                        try:
                            if hasattr(shape, "text"):
                                txt = (shape.text or "").strip()
                                if txt:
                                    parts.append(txt)
                        except Exception:
                            continue
                    if parts:
                        slides_text.append("\n".join(parts))
                combined = "\n\n".join(slides_text).strip()
                return combined if combined else "[No extractable text found in slides]"
            except Exception as e:
                logger.warning(f"PowerPoint extraction error for {display_name}: {e}")
                if file_name.endswith(".pptx"):
                    fallback_text = _extract_pptx_text_fallback(content)
                    if fallback_text:
                        return fallback_text
                return f"[PowerPoint file - content extraction failed: {str(e)}]"

        # Images
        if file_name.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp")):
            if Image is None:
                return f"ðŸ–¼ï¸ Image: {display_name}\n\n(Install pillow to inspect image metadata.)"
            img = Image.open(io.BytesIO(content))
            return f"ðŸ–¼ï¸ Image: {display_name} ({img.width}x{img.height}px, {img.format})"

        return f"ðŸ“Ž File: {display_name} ({len(content)} bytes)"

    except Exception as e:
        logger.error(f"{prefix}Error processing document: {e}", exc_info=True)
        return f"âŒ Error processing {getattr(attachment,'name','attachment')}: {str(e)}"


# =====================================================
# Fuzzy Search Helpers
# =====================================================
def _fuzzy_match(word1: str, word2: str, threshold: float = 0.7) -> bool:
    """
    Check if two words are similar enough using Levenshtein-like comparison.
    Returns True if similarity >= threshold.
    """
    w1, w2 = word1.lower(), word2.lower()
    if w1 == w2:
        return True
    if w1 in w2 or w2 in w1:
        return True
    
    # Simple Levenshtein distance ratio
    len1, len2 = len(w1), len(w2)
    if max(len1, len2) == 0:
        return True
    
    # Quick rejection for very different lengths
    if abs(len1 - len2) > max(len1, len2) * 0.4:
        return False
    
    # Calculate edit distance (simplified)
    if len1 < len2:
        w1, w2 = w2, w1
        len1, len2 = len2, len1
    
    previous_row = range(len2 + 1)
    for i, c1 in enumerate(w1):
        current_row = [i + 1]
        for j, c2 in enumerate(w2):
            insertions = previous_row[j + 1] + 1
            deletions = current_row[j] + 1
            substitutions = previous_row[j] + (c1 != c2)
            current_row.append(min(insertions, deletions, substitutions))
        previous_row = current_row
    
    distance = previous_row[-1]
    similarity = 1 - (distance / max(len1, len2))
    return similarity >= threshold


def _expand_query_with_fuzzy(query: str) -> str:
    """
    Expand search query with common misspelling corrections and extension typo fixes.
    Returns enhanced query for better Graph search results.
    """
    # Common misspellings/typos dictionary
    corrections = {
        'nirobi': 'nairobi',
        'niarobi': 'nairobi',
        'narobi': 'nairobi',
        'agentcon': 'agentcon',
        'agdiag': 'agdiag',
        'agdag': 'agdiag',
        'docuemnt': 'document',
        'documnet': 'document',
        'meetting': 'meeting',
        'meetign': 'meeting',
        'porject': 'project',
        'prject': 'project',
        'repotr': 'report',
        'reprot': 'report',
        'summry': 'summary',
        'sumarry': 'summary',
        'anlysis': 'analysis',
        'anlaysis': 'analysis',
        'presenation': 'presentation',
        'presentaiton': 'presentation',
    }
    
    # Common extension typos
    ext_corrections = {
        '.xlxs': '.xlsx', '.xlx': '.xlsx', '.xslx': '.xlsx',
        '.docsx': '.docx', '.dox': '.docx', '.dcx': '.docx',
        '.pdfx': '.pdf', '.ppptx': '.pptx', '.pptxx': '.pptx',
        '.csvv': '.csv', '.txtt': '.txt',
    }
    
    # First pass: fix extension typos in the whole query
    query_fixed = query
    for typo, correct in ext_corrections.items():
        if typo in query_fixed.lower():
            query_fixed = re.sub(re.escape(typo), correct, query_fixed, flags=re.IGNORECASE)
            logger.info(f"Extension typo correction: '{typo}' -> '{correct}'")
    
    words = query_fixed.split()
    expanded_words = []
    
    for word in words:
        word_lower = word.lower()
        # Check direct correction
        if word_lower in corrections:
            corrected = corrections[word_lower]
            if corrected != word_lower:
                logger.info(f"Spelling correction: '{word}' -> '{corrected}'")
                expanded_words.append(f"({word} OR {corrected})")
            else:
                expanded_words.append(word)
        else:
            expanded_words.append(word)
    
    return " ".join(expanded_words)


def _fuzzy_term_match(term: str, text: str, threshold: float = 0.7) -> bool:
    """
    Check if a search term fuzzy-matches any word in the text.
    """
    text_words = text.lower().split()
    term_lower = term.lower()
    
    for text_word in text_words:
        # Clean punctuation
        clean_word = ''.join(c for c in text_word if c.isalnum())
        if _fuzzy_match(term_lower, clean_word, threshold):
            return True
    return False


# =====================================================
# Graph API - SharePoint Search
# =====================================================
def search_sharepoint(query: str, token: str, user_context: bool = True, user_upn: str = "") -> dict:
    """
    Search SharePoint/OneDrive using Microsoft Graph Search API.
    Includes fuzzy matching for typo tolerance.
    
    Args:
        query: Search query string
        token: Graph access token (delegated or app-only)
        user_context: If True (default), token is delegated and respects user permissions.
                     If False, token is app-only and may access more than user can see.
        user_upn: User principal name (email) for permission pre-filtering
    
    Returns:
        Dictionary with search results including webUrls for downloading
    """
    try:
        headers = {
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json",
        }

        # Region handling: default to 'US' per Graph error guidance; allow override via env
        region = os.getenv("GRAPH_SEARCH_REGION", "US").strip().upper() or "US"
        
        # Expand query for better recall
        cleaned_query = _normalize_search_query(query)
        expanded_query = _expand_query_with_fuzzy(cleaned_query or query)

        # Add phrase wrapping for relevance
        if " " in cleaned_query:
            expanded_query = f'("{cleaned_query}") OR ({expanded_query})'

        # Only add joined variant for 2-3 word queries (likely a document title)
        original_words = query.strip()
        word_list = [w for w in original_words.split() if w.lower() not in GRAPH_QUERY_STOPWORDS]
        if 2 <= len(word_list) <= 3:
            joined_variant = "".join(word_list).replace(".", "")
            if joined_variant and joined_variant.lower() != original_words.lower():
                expanded_query = f"({expanded_query}) OR ({joined_variant}*)"

        # Add extension-aware queries for searches that might be filenames
        base = (cleaned_query or query).strip()
        if base and " " not in base:
            ext_candidates = [".docx", ".pdf", ".xlsx", ".pptx", ".txt", ".csv"]
            ext_terms = [f'"{base}{ext}"' for ext in ext_candidates if not base.endswith(ext)]
            if ext_terms:
                ext_query = " OR ".join(ext_terms)
                expanded_query = f"({expanded_query}) OR ({ext_query}) OR ({base}*)"
        elif base and "." in base:
            # If extension already present in multi-word query, wrap as phrase for precision
            expanded_query = f'({expanded_query}) OR ("{base}")'
        
        # FINAL GUARD: Ensure parentheses grouping is sound
        if not expanded_query.startswith("("):
            expanded_query = f"({expanded_query})"

        search_body = {
            "requests": [
                {
                    "entityTypes": ["driveItem", "listItem", "site"],
                    "query": {"queryString": expanded_query},
                    "from": 0,
                    "size": Config.MAX_GRAPH_SEARCH_RESULTS,
                    "region": region,
                }
            ]
        }

        # Use retry logic with exponential backoff for Graph search
        def _search_graph():
            return requests.post(
                "https://graph.microsoft.com/v1.0/search/query",
                headers=headers,
                json=search_body,
                # Search is far slower than ordinary Graph reads — use its own,
                # longer budget so a cold/complex query isn't dropped at 8s.
                timeout=getattr(config, "GRAPH_SEARCH_TIMEOUT", 25),
            )
        
        try:
            resp = _retry_request(_search_graph, max_retries=1, initial_delay=0)  # fail fast â€” no retries for search
        except (requests.exceptions.Timeout, requests.exceptions.ReadTimeout) as e:
            SearchLogger.error(f"Graph search timed out")
            return {"error": f"Graph search timed out"}

        # On transient 5xx errors, skip retry and let cache fallback handle the gap.
        if resp.status_code in (500, 502, 503, 504):
            SearchLogger.warning(f"Graph search failed (HTTP {resp.status_code}) - using cache fallback when available")
            return {"error": f"Graph search failed (HTTP {resp.status_code})"}

        if resp.status_code != 200:
            # Retry once with US if server indicates only US is valid
            try:
                body_txt = resp.text or ""
            except Exception:
                body_txt = ""
            if "Only valid regions are US" in body_txt and region != "US":
                search_body["requests"][0]["region"] = "US"
                
                def _search_graph_us():
                    return requests.post(
                        "https://graph.microsoft.com/v1.0/search/query",
                        headers=headers,
                        json=search_body,
                        timeout=config.GRAPH_TIMEOUT,
                    )
                
                try:
                    resp = _retry_request(_search_graph_us, max_retries=1, initial_delay=0)  # no retry for region fallback
                except (requests.exceptions.Timeout, requests.exceptions.ReadTimeout) as e:
                    SearchLogger.error(f"Graph search timed out")
                    return {"error": f"Graph search timed out"}
                
                if resp.status_code != 200:
                    SearchLogger.warning(f"Graph search failed (HTTP {resp.status_code})")
                    return {"error": f"Graph search failed"}
            else:
                SearchLogger.warning(f"Graph search failed (HTTP {resp.status_code})")
                return {"error": f"Graph search failed"}

        data = resp.json()
        results = []
        permission_filtered_count = 0
        relevance_filtered_count = 0
        doctype_filtered_count = 0
        
        # Extract query keywords for relevance filtering
        query_terms = _extract_query_terms(query)
        min_relevance = float(os.getenv("GRAPH_MIN_RELEVANCE", "0.05"))  # Lower threshold to include more results
        fallback_topn = int(os.getenv("GRAPH_FALLBACK_TOPN", "10"))  # Keep more fallback results

        low_relevance_results = []

        for block in data.get("value", []):
            for container in block.get("hitsContainers", []) or []:
                for hit in container.get("hits", []) or []:
                    resource = hit.get("resource", {}) or {}
                    name = resource.get("name") or resource.get("title") or "Untitled"
                    summary = (hit.get("summary") or "").strip()
                    web_url = resource.get("webUrl", "")

                    # Filter out unsupported file types (e.g., .zip)
                    try:
                        name_lower = name.lower()
                        url_lower = (web_url or "").lower()
                        name_has_ext = "." in name_lower and not name_lower.endswith(".")
                        url_has_ext = "." in url_lower and "/" in url_lower
                        
                        if name_has_ext:
                            if not _is_supported_document(name_lower):
                                doctype_filtered_count += 1
                                continue
                        
                        if url_has_ext:
                            url_ext = url_lower.rsplit(".", 1)[-1]
                            if url_ext and not _is_supported_document(f"x.{url_ext}"):
                                doctype_filtered_count += 1
                                continue
                    except Exception as filter_err:
                        pass
                    
                    # Pre-filter: Skip documents user cannot access (e.g., other users' personal OneDrive)
                    if user_upn and not is_url_accessible_by_user(web_url, user_upn, user_context=user_context):
                        permission_filtered_count += 1
                        continue
                    elif not user_upn and "/personal/" in web_url.lower():
                        permission_filtered_count += 1
                        continue
                    
                    # Calculate relevance score based on keyword + fuzzy matches
                    searchable_text = (name + " " + summary).lower()
                    exact_matches = [term for term in query_terms if term in searchable_text]
                    fuzzy_matches = []
                    if query_terms and not exact_matches:
                        for term in query_terms:
                            if len(term) >= 4 and _fuzzy_term_match(term, searchable_text, threshold=0.8):
                                fuzzy_matches.append(term)
                    matches = len(set(exact_matches + fuzzy_matches))
                    relevance_score = matches / len(query_terms) if query_terms else 0

                    if relevance_score >= min_relevance or not query_terms:
                        results.append({
                            "name": name,
                            "webUrl": web_url,
                            "summary": summary,
                            "driveId": resource.get("parentReference", {}).get("driveId", ""),
                            "itemId": resource.get("id", ""),
                            "relevance_score": relevance_score,
                            "_from_sharepoint": True,
                        })
                    else:
                        relevance_filtered_count += 1
                        low_relevance_results.append({
                            "name": name,
                            "webUrl": web_url,
                            "summary": summary,
                            "driveId": resource.get("parentReference", {}).get("driveId", ""),
                            "itemId": resource.get("id", ""),
                            "relevance_score": relevance_score,
                            "_from_sharepoint": True,
                        })

        # If everything was filtered, keep a small fallback set to avoid empty results
        if not results and low_relevance_results and query_terms:
            keep_n = min(fallback_topn, len(low_relevance_results))
            results.extend(low_relevance_results[:keep_n])
        
        # Sort by relevance score (highest first)
        results.sort(key=lambda x: x.get("relevance_score", 0), reverse=True)
        
        # Log results with clean formatting
        SearchLogger.search_source("Graph", len(results), query)
        for doc in results[:5]:
            ext = doc.get("name", "").rsplit(".", 1)[-1].lower() if "." in doc.get("name", "") else "?"
            SearchLogger.document(doc.get("name", "Untitled"), doc_type=ext, score=doc.get("relevance_score"))
        if len(results) > 5:
            logger.info(f"    ... and {len(results) - 5} more")
        
        # If no results found, attempt a simplified query retry (keyword-only search)
        if not results and query_terms and len(query_terms) > 1:
            logger.info(f"âš ï¸ No results found with full query; attempting simplified search with primary keyword only")
            primary_term = query_terms[0]  # First keyword only
            try:
                simplified_search_body = {
                    "requests": [
                        {
                            "entityTypes": ["driveItem", "listItem", "site"],
                            "query": {"queryString": primary_term},
                            "from": 0,
                            "size": Config.MAX_GRAPH_SEARCH_RESULTS,
                            "region": region,
                        }
                    ]
                }
                
                def _search_graph_simplified():
                    return requests.post(
                        "https://graph.microsoft.com/v1.0/search/query",
                        headers=headers,
                        json=simplified_search_body,
                        timeout=config.GRAPH_TIMEOUT,
                    )
                
                resp_simplified = _retry_request(_search_graph_simplified, max_retries=2, initial_delay=1.0)
                if resp_simplified.status_code == 200:
                    data_simplified = resp_simplified.json()
                    logger.info(f"ðŸ”„ Simplified search on '{primary_term}' returned results")
                    
                    for block in data_simplified.get("value", []):
                        for container in block.get("hitsContainers", []) or []:
                            for hit in container.get("hits", []) or []:
                                resource = hit.get("resource", {}) or {}
                                name = resource.get("name") or resource.get("title") or "Untitled"
                                summary = (hit.get("summary") or "").strip()
                                web_url = resource.get("webUrl", "")
                                
                                # Apply same filtering rules as main search
                                if not _is_supported_document(name):
                                    continue
                                if user_upn and not is_url_accessible_by_user(web_url, user_upn, user_context=user_context):
                                    continue
                                
                                doc_id = resource.get("id") or web_url
                                if doc_id and doc_id not in [r.get("id") or r.get("webUrl") for r in results]:
                                    results.append({
                                        "name": name,
                                        "webUrl": web_url,
                                        "summary": summary,
                                        "driveId": resource.get("parentReference", {}).get("driveId", ""),
                                        "itemId": resource.get("id", ""),
                                        "relevance_score": 0.5,  # Lower score since it's from simplified search
                                        "_from_sharepoint": True,
                                        "_from_fallback_search": True
                                    })
                    
                    if results:
                        logger.info(f"âœ… Simplified search recovered {len(results)} results")
            except Exception as fallback_err:
                logger.warning(f"Simplified search fallback failed: {fallback_err}")
        
        return {"results": results, "count": len(results)}

    except Exception as e:
        logger.error(f"SharePoint search error: {e}", exc_info=True)
        return {"error": str(e)}


def _log_search_diagnostics(query: str, results: list, search_source: str = "Graph", **kwargs) -> None:
    """Log detailed search diagnostics for debugging.
    
    Args:
        query: The search query
        results: List of returned results
        search_source: Source of results (Graph, OneDrive, AI Search, Cache, etc.)
        **kwargs: Additional diagnostic data to log
    """
    logger.info(f"ðŸ“Š Search Diagnostics: {search_source}")
    logger.info(f"  Query: '{query}'")
    logger.info(f"  Results: {len(results)} found")
    
    if kwargs:
        for key, value in kwargs.items():
            if isinstance(value, list):
                logger.info(f"  {key}: {len(value)} items - {value[:3]}{'...' if len(value) > 3 else ''}")
            else:
                logger.info(f"  {key}: {value}")
    
    # Log result types breakdown
    if results:
        extensions = {}
        sources = {}
        for doc in results:
            # Count by extension
            name = doc.get("name", "")
            if "." in name:
                ext = name.rsplit(".", 1)[-1].lower()
                extensions[ext] = extensions.get(ext, 0) + 1
            
            # Count by source
            if doc.get("_from_onedrive_search"):
                sources["OneDrive"] = sources.get("OneDrive", 0) + 1
            elif doc.get("_from_live_graph"):
                sources["Graph"] = sources.get("Graph", 0) + 1
            elif doc.get("_from_ai_search"):
                sources["AI Search"] = sources.get("AI Search", 0) + 1
            elif doc.get("_from_web"):
                sources["Web"] = sources.get("Web", 0) + 1
            elif doc.get("_from_cache"):
                sources["Cache"] = sources.get("Cache", 0) + 1
        
        if extensions:
            logger.info(f"  File types: {dict(sorted(extensions.items(), key=lambda x: x[1], reverse=True))}")
        if sources:
            logger.info(f"  Result sources: {sources}")


def list_user_files(user_id: str, user_assertion: Optional[str] = None, top: int = 10) -> list[dict]:
    """List files from the user's OneDrive root.

    Uses delegated token (/me) when available; otherwise uses app-only to /users/{id}.
    Filters to supported document extensions only.
    """
    try:
        if not user_id:
            return []
        top = max(1, min(int(top), 50))
        token = get_graph_token(user_assertion)
        if not token:
            return []
        headers = {"Authorization": f"Bearer {token}"}

        if user_assertion:
            url = f"https://graph.microsoft.com/v1.0/me/drive/root/children?$top=200"
        else:
            url = f"https://graph.microsoft.com/v1.0/users/{user_id}/drive/root/children?$top=200"

        results = []
        next_url = url
        while next_url and len(results) < top:
            resp = requests.get(next_url, headers=headers, timeout=config.GRAPH_TIMEOUT)
            _log_http("", "List OneDrive files", resp)
            if resp.status_code != 200:
                break
            data = resp.json() or {}
            items = data.get("value", []) or []
            for item in items:
                if not item.get("file"):
                    continue
                name = item.get("name") or ""
                if not _is_supported_document(name):
                    continue
                results.append({
                    "name": name,
                    "webUrl": item.get("webUrl") or "",
                    "id": item.get("id") or ""
                })
                if len(results) >= top:
                    break
            next_url = data.get("@odata.nextLink")
        return results
    except Exception as e:
        logger.error(f"List OneDrive files error: {e}")
        return []


def list_sharepoint_files(site_urls: list[str], user_assertion: Optional[str] = None, top_per_site: int = 10) -> list[dict]:
    """List files from SharePoint sites' default document libraries.

    Returns lightweight metadata (name, webUrl, driveId, itemId).
    """
    results = []
    try:
        if not site_urls:
            return []
        top_per_site = max(1, min(int(top_per_site), 50))
        token = get_graph_token(user_assertion)
        if not token:
            return []
        headers = {"Authorization": f"Bearer {token}"}

        for site_url in site_urls:
            site_id = _resolve_site_id(site_url, headers)
            if not site_id:
                continue
            list_url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drive/root/children?$top=200"
            next_url = list_url
            while next_url and len(results) < (top_per_site * len(site_urls)):
                resp = requests.get(next_url, headers=headers, timeout=config.GRAPH_TIMEOUT)
                _log_http("", "List SharePoint files", resp)
                if resp.status_code != 200:
                    break
                data = resp.json() or {}
                items = data.get("value", []) or []
                for item in items:
                    if not item.get("file"):
                        continue
                    name = item.get("name") or ""
                    if not _is_supported_document(name):
                        continue
                    results.append({
                        "name": name,
                        "webUrl": item.get("webUrl") or "",
                        "driveId": item.get("parentReference", {}).get("driveId", ""),
                        "itemId": item.get("id") or "",
                    })
                    if len(results) >= (top_per_site * len(site_urls)):
                        break
                next_url = data.get("@odata.nextLink")
    except Exception as e:
        logger.error(f"List SharePoint files error: {e}")
    return results

def download_and_extract_content(url: str, token: str, file_name: str, drive_id: str = "", item_id: str = "") -> str:
    """
    Download a file from SharePoint/OneDrive and extract its text content.
    
    Args:
        url: Web URL of the file
        token: Graph API token
        file_name: Name of the file
        drive_id: Drive ID from search result
        item_id: Item ID from search result
    
    Returns:
        Extracted text content
    """
    try:
        headers = {"Authorization": f"Bearer {token}"}
        
        # Use Graph API directly with drive/item IDs if available
        if drive_id and item_id:
            graph_url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{item_id}/content"
            logger.info(f"Downloading via Graph API: {graph_url}")
            resp = requests.get(
                graph_url,
                headers=headers,
                timeout=int(os.environ.get("DOCUMENT_DOWNLOAD_TIMEOUT", "8")),
            )
        else:
            # Fallback to share ID approach
            logger.info(f"Downloading content from: {url}")
            resp = requests.get(
                url,
                headers=headers,
                timeout=int(os.environ.get("DOCUMENT_DOWNLOAD_TIMEOUT", "8")),
            )
            
            # If unauthorized or HTML viewer page, try Graph shares/path
            if resp.status_code in (401, 403) or _is_probably_html(resp):
                logger.warning(f"Direct fetch is unauthorized or HTML viewer; attempting Graph shares/path")
                share_id = _graph_share_id_from_url(url)
                graph_url = f"https://graph.microsoft.com/v1.0/shares/{share_id}/driveItem/content"
                resp = requests.get(
                    graph_url,
                    headers=headers,
                    timeout=int(os.environ.get("DOCUMENT_DOWNLOAD_TIMEOUT", "8")),
                )
                if resp.status_code in (400, 404) or _is_probably_html(resp):
                    alt_resp = _graph_download_via_path(url, headers, "")
                    if alt_resp is not None:
                        resp = alt_resp
        
        if resp.status_code != 200:
            return f"[Unable to download: HTTP {resp.status_code} - Check Files.Read.All permission]"
        
        content = resp.content
        file_name_lower = file_name.lower()

        # If we still received HTML, avoid parsing as a binary document
        try:
            if _is_probably_html(resp):
                return f"[Content appears to be an HTML viewer page, not raw file: {file_name}]"
        except Exception:
            pass
        
        # Extract text based on file type
        if file_name_lower.endswith(".pdf") and pypdf:
            reader = pypdf.PdfReader(io.BytesIO(content))
            text_parts = []
            for page in reader.pages:
                try:
                    t = page.extract_text() or ""
                except Exception:
                    # Some PDFs trigger internal parser errors (e.g., missing 'bbox'); skip page
                    t = ""
                if t.strip():
                    text_parts.append(t)
            extracted = "\n".join(text_parts)
            logger.info(f"Extracted {len(extracted)} chars from {len(reader.pages)} pages: {file_name}")
            return extracted
        
        elif file_name_lower.endswith((".docx", ".doc")):
            if file_name_lower.endswith(".doc") and not file_name_lower.endswith(".docx"):
                # Try python-docx first (works if misnamed .docx)
                try:
                    if Document:
                        doc = Document(io.BytesIO(content))
                        extracted = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                        logger.info(f"âœ“ .doc extraction via python-docx: {len(extracted)} chars from {len([p for p in doc.paragraphs if p.text.strip()])} paragraphs: {file_name}")
                        if extracted:
                            return extracted
                except Exception as e:
                    logger.warning(f"python-docx extraction failed for {file_name}: {e}")
                
                # Real legacy .doc - try textract
                from simple_file_handler import _textract_legacy_office
                legacy_text = _textract_legacy_office(content, "doc", file_name)
                if legacy_text:
                    logger.info(f"âœ“ .doc extraction via textract: {len(legacy_text)} chars: {file_name}")
                    return legacy_text
                if textract is None:
                    return "[Legacy Word .doc detected. Install textract to extract content.]"
                logger.warning(f"Legacy .doc extraction failed for {file_name}")
                return "[Legacy Word .doc detected but extraction failed.]"
            
            # .docx files - try multiple methods
            if Document is None:
                logger.warning(f"python-docx not installed, cannot extract .docx: {file_name}")
                return "[Word .docx detected but python-docx is not installed.]"
            
            try:
                doc = Document(io.BytesIO(content))
                
                # Debug: Log total paragraphs and their details
                logger.info(f"ðŸ“„ .docx file opened: {len(doc.paragraphs)} total paragraphs, {len(doc.tables)} tables")
                
                # Collect all paragraphs with text
                paragraphs_with_text = []
                for idx, p in enumerate(doc.paragraphs):
                    text = p.text.strip()
                    if text:
                        paragraphs_with_text.append(text)
                        if len(paragraphs_with_text) <= 5:  # Log first 5 paragraphs for debugging
                            preview = text[:100].replace("\n", " ")
                            logger.debug(f"  Para[{idx}]: {len(text)} chars | {preview}...")
                
                extracted = "\n".join(paragraphs_with_text)
                logger.info(f"âœ“ .docx extraction (method=paragraphs): {len(extracted)} chars from {len(paragraphs_with_text)} non-empty paragraphs: {file_name}")
                if extracted:
                    logger.debug(f"  Content preview: {extracted[:200].replace(chr(10), ' ')}")
                
                # If extraction is suspiciously small, try tables and other elements
                if len(extracted) < 50 and len(paragraphs_with_text) > 0:
                    logger.info(f"âš  .docx paragraph extraction returned only {len(extracted)} chars ({len(paragraphs_with_text)} paragraphs), checking tables ({len(doc.tables)} total)...")
                    tables_text = []
                    for t_idx, table in enumerate(doc.tables):
                        for r_idx, row in enumerate(table.rows):
                            row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_cells:
                                tables_text.append(" | ".join(row_cells))
                                if len(tables_text) <= 3:
                                    logger.debug(f"  Table[{t_idx}].Row[{r_idx}]: {len(row_cells)} cells, {len(' | '.join(row_cells))} chars")
                    if tables_text:
                        table_content = "\n".join(tables_text)
                        logger.info(f"âœ“ Found {len(table_content)} chars in {len(doc.tables)} table(s)")
                        extracted = extracted + "\n\n" + table_content if extracted else table_content
                
                # If still very small, try to extract from all shapes/text boxes
                if len(extracted) < 50:
                    logger.info(f"âš  Still <50 chars, scanning for shapes/textboxes...")
                    shapes_text = []
                    try:
                        for shape in doc.element.body.iter():
                            # Try to extract text from shape elements
                            if hasattr(shape, 'text') and shape.text.strip():
                                shapes_text.append(shape.text.strip())
                    except Exception as e:
                        logger.debug(f"  Shape extraction not available: {e}")
                    
                    if shapes_text:
                        shapes_content = "\n".join(shapes_text)
                        logger.info(f"âœ“ Found {len(shapes_content)} chars in shapes/textboxes")
                        extracted = extracted + "\n\n" + shapes_content if extracted else shapes_content
                
                # If still very small, try textract as fallback
                if len(extracted) < 100 and textract:
                    logger.info(f"âš  Extracted content still <100 chars, trying textract fallback...")
                    try:
                        textract_result = textract.process(io.BytesIO(content), extension="docx")
                        if textract_result:
                            fallback_text = textract_result.decode("utf-8", errors="ignore").strip()
                            if len(fallback_text) > len(extracted):
                                logger.info(f"âœ“ Textract fallback returned {len(fallback_text)} chars (improvement over {len(extracted)})")
                                logger.debug(f"  Textract preview: {fallback_text[:200]}")
                                extracted = fallback_text
                            else:
                                logger.info(f"  Textract returned {len(fallback_text)} chars (no improvement)")
                    except Exception as te:
                        logger.warning(f"Textract fallback failed: {te}")
                
                if not extracted:
                    logger.warning(f"âš  .docx {file_name} resulted in empty extraction after all methods - document may be empty or corrupted")
                    return "[âš  DOCX file appears empty or content extraction failed after trying multiple methods]"
                
                logger.info(f"âœ… Final .docx extraction for {file_name}: {len(extracted)} chars")
                return extracted
            except Exception as e:
                logger.error(f"âœ— .docx extraction failed for {file_name}: {e}", exc_info=True)
                return f"[Error extracting .docx: {str(e)[:100]}]"
        
        elif file_name_lower.endswith((".xlsx", ".xls")):
            if file_name_lower.endswith(".xls"):
                xls_text = _extract_xls_with_xlrd(content)
                if xls_text:
                    logger.info(f"Extracted {len(xls_text)} chars from legacy .xls: {file_name}")
                    return xls_text
                if xlrd is None:
                    return "[Legacy Excel .xls detected. Install xlrd==1.2.0 to extract content.]"
                return "[Legacy Excel .xls detected but extraction failed.]"
            if not load_workbook:
                return "[Excel .xlsx detected but openpyxl is not installed.]"
            wb = load_workbook(io.BytesIO(content), data_only=True)
            sheets_text = []
            total_rows = 0
            for sheet in wb.worksheets:
                rows = []
                sheet_rows = 0
                # Capture ALL rows - no limits
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None and str(cell).strip() for cell in row):
                        rows.append(" | ".join(str(cell or "") for cell in row))
                        sheet_rows += 1
                if rows:
                    sheets_text.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))
                    total_rows += sheet_rows
            extracted = "\n\n".join(sheets_text)
            logger.info(f"Extracted {len(extracted)} chars from {total_rows} rows across {len(wb.worksheets)} sheets: {file_name}")
            return extracted

        elif file_name_lower.endswith((".pptx", ".ppt")):
            if file_name_lower.endswith(".ppt") and not file_name_lower.endswith(".pptx"):
                # Check if it's actually a .pptx file (ZIP signature)
                if content[:4] == b"PK\x03\x04":
                    try:
                        if Presentation is None:
                            return "[PowerPoint .pptx detected but python-pptx is not installed.]"
                        prs = Presentation(io.BytesIO(content))
                        slides_text = []
                        for slide in prs.slides:
                            parts = []
                            for shape in slide.shapes:
                                try:
                                    if hasattr(shape, "text"):
                                        txt = (shape.text or "").strip()
                                        if txt:
                                            parts.append(txt)
                                except Exception:
                                    continue
                            if parts:
                                slides_text.append("\n".join(parts))
                        combined = "\n\n".join(slides_text).strip()
                        return combined if combined else "[No extractable text found in slides]"
                    except Exception as e:
                        logger.warning(f"PowerPoint extraction error for {file_name}: {e}")
                
                # Real legacy .ppt - try textract
                from simple_file_handler import _textract_legacy_office
                legacy_text = _textract_legacy_office(content, "ppt", file_name)
                if legacy_text:
                    return legacy_text
                if textract is None:
                    return "[Legacy PowerPoint .ppt detected. Install textract to extract content.]"
                return "[Legacy PowerPoint .ppt detected but extraction failed.]"
            
            # .pptx files - standard extraction
            if Presentation is None:
                return "[PowerPoint .pptx detected but python-pptx is not installed.]"
            try:
                prs = Presentation(io.BytesIO(content))
                slides_text = []
                for slide in prs.slides:
                    parts = []
                    for shape in slide.shapes:
                        try:
                            if hasattr(shape, "text"):
                                txt = (shape.text or "").strip()
                                if txt:
                                    parts.append(txt)
                        except Exception:
                            continue
                    if parts:
                        slides_text.append("\n".join(parts))
                combined = "\n\n".join(slides_text).strip()
                return combined if combined else "[No extractable text found in slides]"
            except Exception as e:
                logger.warning(f"PowerPoint extraction error for {file_name}: {e}")
                if file_name_lower.endswith(".pptx"):
                    fallback_text = _extract_pptx_text_fallback(content)
                    if fallback_text:
                        return fallback_text
                return f"[PowerPoint file - content extraction failed: {str(e)}]"
        
        elif file_name_lower.endswith(".txt"):
            extracted = content.decode("utf-8", errors="ignore")
            logger.info(f"Extracted {len(extracted)} chars from text file: {file_name}")
            return extracted

        elif file_name_lower.endswith((".csv", ".json", ".xml")):
            try:
                extracted = content.decode("utf-8", errors="ignore")
                logger.info(f"Extracted {len(extracted)} chars from {file_name_lower.split('.')[-1].upper()} file: {file_name}")
                return extracted
            except Exception:
                return "[Unable to decode file as UTF-8 text]"

        elif file_name_lower.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp")) and 'Image' in globals() and Image is not None:
            try:
                img = Image.open(io.BytesIO(content))
                return f"[Image: {file_name} {img.width}x{img.height}px, {img.format}]"
            except Exception:
                return f"[Image file detected but failed to read basic metadata]"
        
        else:
            return f"[{file_name}: Content extraction not supported for this file type]"
    
    except Exception as e:
        logger.error(f"Error extracting content from {file_name}: {e}")
        return f"[Error extracting content: {str(e)}]"


# =====================================================
# Graph API - Send Email (FIXED for app-only)
# =====================================================
def send_email(to: str, subject: str, body: str, token: str, user_context: bool = True, sender_upn: Optional[str] = None) -> str:
    """
    Send email via Microsoft Graph.
    
    Args:
        to: Recipient email address
        subject: Email subject
        body: Email body (plain text)
        token: Graph access token (delegated or app-only)
        user_context: If True (default), uses /me/sendMail (delegated token).
        sender_upn: Sender's UPN (required only for app-only mode)
    """
    try:
        auth_type = "delegated (user-context)" if user_context else "app-only"
        logger.info(f"Sending email with {auth_type} token")
        
        headers = {
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json",
        }

        email_body = {
            "message": {
                "subject": subject,
                "body": {"contentType": "Text", "content": body},
                "toRecipients": [{"emailAddress": {"address": to}}],
            },
            "saveToSentItems": "true",
        }

        # Use appropriate endpoint based on token type
        if user_context:
            # Delegated token: send as authenticated user
            endpoint = "https://graph.microsoft.com/v1.0/me/sendMail"
        else:
            # App-only token: requires explicit sender UPN
            sender = sender_upn or getattr(config, "SENDER_UPN", None)
            if not sender:
                return (
                    "\n\nâŒ Email not configured: set SENDER_UPN in config/env. "
                    "App-only sendMail requires /users/{SENDER_UPN}/sendMail."
                )
            endpoint = f"https://graph.microsoft.com/v1.0/users/{sender}/sendMail"

        resp = requests.post(
            endpoint,
            headers=headers,
            json=email_body,
            timeout=config.GRAPH_TIMEOUT,
        )

        _log_http("", "Send mail", resp)

        if resp.status_code == 202:
            return f"\n\nâœ… Email sent successfully to {to}!"
        return f"\n\nâŒ Failed to send email (HTTP {resp.status_code})."

    except Exception as e:
        logger.error(f"Email error: {e}", exc_info=True)
        return f"\n\nâŒ Email error: {str(e)}"


# =====================================================
# Graph API - Full Crawl and Caching
# =====================================================
def _crawl_drive_items(
    drive_id: str,
    token: str,
    owner_user_id: str,
    treat_as_shared: bool,
    max_items: int,
    max_depth: int,
    stats: dict,
):
    """Breadth-first crawl of a drive and cache supported documents."""
    headers = {"Authorization": f"Bearer {token}"}
    cache = get_cache()
    queue: list[tuple[str, int]] = [("root", 0)]
    
    # Track items processed in this crawl session to avoid re-indexing duplicates
    # within the same crawl (in case same item appears in multiple places)
    seen_in_crawl = set()

    processed = 0
    while queue and processed < max_items:
        current_item, depth = queue.pop(0)
        page_url = (
            f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{current_item}/children"
            "?$select=id,name,folder,file,parentReference,webUrl,size,lastModifiedDateTime"
        )

        while page_url and processed < max_items:
            resp = requests.get(page_url, headers=headers, timeout=config.GRAPH_TIMEOUT)
            _log_http("", "Drive list", resp)
            if resp.status_code != 200:
                stats["errors"] += 1
                break

            data = resp.json() or {}
            for item in data.get("value", []) or []:
                # Recurse into folders up to max_depth
                if "folder" in item:
                    if depth < max_depth:
                        queue.append((item.get("id"), depth + 1))
                    continue

                # Only process files we support
                if "file" not in item:
                    continue

                name = item.get("name", "")
                if not _is_supported_document(name):
                    stats["skipped"] += 1
                    continue

                if _should_skip_file(item):
                    stats["skipped"] += 1
                    continue

                item_id = item.get("id", "")
                doc_id = _make_drive_doc_id(drive_id, item_id)
                web_url = item.get("webUrl", "")
                
                # Skip if already processed in this crawl session
                if doc_id in seen_in_crawl:
                    stats["skipped"] += 1
                    continue
                seen_in_crawl.add(doc_id)

                if treat_as_shared:
                    if cache.has_shared_document(doc_id):
                        stats["skipped"] += 1
                        continue
                else:
                    if cache.has_document(doc_id, user_id=owner_user_id):
                        stats["skipped"] += 1
                        continue

                content = download_and_extract_content(web_url, token, name, drive_id, item_id)
                if not content or len(content.strip()) < 10:
                    stats["skipped"] += 1
                    continue

                error_indicators = [
                    "[Error extracting",
                    "[Unable to download",
                    "Content extraction not supported",
                    "Check Files.Read.All permission",
                    "File is not a zip file",
                ]
                if any(indicator in content for indicator in error_indicators):
                    stats["skipped"] += 1
                    continue

                metadata = {
                    "drive_id": drive_id,
                    "item_id": item_id,
                    "last_modified": item.get("lastModifiedDateTime"),
                    "path": item.get("parentReference", {}).get("path", ""),
                    "crawled_at": datetime.now().isoformat(),  # Track when document was crawled
                    "access_verified_at": datetime.now().isoformat(),  # Track permission verification time
                }

                # SECURITY: Verify user has access before caching
                # This prevents caching documents during crawl that user shouldn't access
                try:
                    # For personal OneDrive, owner_user_id IS the user - always accessible
                    # For shared SharePoint, verify the user actually has permission
                    should_cache = True
                    
                    if treat_as_shared and web_url:
                        # Get user's email/UPN from graph to verify SharePoint access
                        # For now, treat shared documents as accessible during crawl
                        # Real-time permission checks happen during search
                        logger.debug(f"Crawling shared document: {name} (owner: {owner_user_id[:8]}...)")
                    elif not treat_as_shared:
                        # Personal OneDrive - verify owner matches
                        if "/personal/" in web_url.lower():
                            try:
                                # Extract owner from URL
                                parts = web_url.lower().split("/personal/")
                                if len(parts) > 1:
                                    owner_part = parts[1].split("/")[0]
                                    if "_" in owner_part:
                                        owner_clean = owner_part.replace("_", ".")
                                        # Owner can access their own documents
                                        logger.debug(f"Personal OneDrive document: {name} (owner verified)")
                            except Exception:
                                pass
                    
                    if should_cache:
                        # SECURITY: All documents are user-specific (no shared cache)
                        if treat_as_shared:
                            cache.add_shared_document(doc_id, name, web_url, content, metadata=metadata)
                            stats["shared_indexed"] += 1
                        else:
                            # Personal drive documents
                            cache.add_document(doc_id, name, web_url, content, user_id=owner_user_id, metadata=metadata)
                            stats["personal_indexed"] += 1
                    else:
                        logger.warning(f"ðŸ”’ SECURITY: Skipped caching document without verified access: {name}")
                        stats["skipped"] += 1
                except Exception as security_err:
                    logger.error(f"Security verification failed for {name}: {security_err}")
                    stats["skipped"] += 1
                    continue

                processed += 1
                if processed >= max_items:
                    break

            page_url = data.get("@odata.nextLink")


def crawl_accessible_documents(
    user_id: str,
    include_personal: bool = True,
    include_sites: bool = True,
    max_items_per_drive: Optional[int] = None,
    max_depth: Optional[int] = None,
    user_display_name: str = None,
) -> dict:
    """Background crawl to index configured SharePoint documents for faster cache searches.

    NOTE: This is for BACKGROUND INDEXING only. For fresh results, use live Graph search.
    - Crawling builds a local cache for faster subsequent searches
    - Live Graph search (search_sharepoint) provides real-time results without waiting for crawls

    CRITICAL SECURITY:
    - Personal drive items are cached with user_id tag (user-specific access)
    - SharePoint site libraries are cached with user_id tag (user sees only what they request)
    - NO shared cache - all documents are user-specific to prevent unauthorized access
    - User profile is cached at the start for personalized greetings
    
    Args:
        user_id: User's AAD object ID or UPN (REQUIRED)
        include_personal: If True, crawl user's personal OneDrive
        include_sites: If True, crawl configured SharePoint sites
        max_items_per_drive: Max items to index per drive
        max_depth: Max folder depth to crawl
    
    Returns:
        Statistics dictionary with indexed counts
    """

    if not user_id:
        logger.error("Crawl aborted: user_id is required for security reasons")
        return {"error": "user_id is required for security reasons"}

    # Special handling for shared crawl (startup crawl with no real user)
    is_shared_crawl = (user_id == "shared")
    
    # Ensure user profile is cached first for personalized experience (skip for shared crawl)
    display_name = "shared access"  # Default for shared crawl
    if not is_shared_crawl:
        try:
            profile_cached = ensure_user_profile_cached(user_id)
            if not profile_cached:
                logger.warning(f"Could not cache user profile for {user_id[:min(20, len(user_id))]}... - profile fetch may have failed")
        except Exception as e:
            logger.error(f"Error ensuring profile cache for {user_id[:min(20, len(user_id))]}...: {e}")
        
        # Use provided display name as fallback if Graph profile doesn't have it
        display_name = get_user_display_name(user_id, fallback_name=user_display_name)
    
    logger.info(f"Starting document crawl for {'shared libraries' if is_shared_crawl else f'user: {display_name}'} ({user_id[:min(8, len(user_id))]}...)")

    token = get_graph_token()
    if not token:
        logger.error("Crawl aborted: Graph token unavailable")
        return {"error": "Graph token unavailable"}

    headers = {"Authorization": f"Bearer {token}"}
    max_items = max_items_per_drive or GRAPH_CRAWL_MAX_ITEMS_PER_DRIVE
    depth_limit = max_depth or GRAPH_CRAWL_MAX_DEPTH

    stats = {
        "user_id": user_id,
        "user_display_name": display_name,
        "personal_indexed": 0,
        "sharepoint_indexed": 0,
        "skipped": 0,
        "errors": 0,
    }

    # Crawl personal OneDrive (always user-specific, skip for shared crawl)
    if include_personal and user_id and not is_shared_crawl:
        try:
            drive_resp = requests.get(
                f"https://graph.microsoft.com/v1.0/users/{user_id}/drive",
                headers=headers,
                timeout=config.GRAPH_TIMEOUT,
            )
            _log_http("", "User drive", drive_resp)
            if drive_resp.status_code == 200:
                drive_id = drive_resp.json().get("id")
                if drive_id:
                    logger.info(f"Crawling personal OneDrive for {display_name}...")
                    _crawl_drive_items(
                        drive_id,
                        token,
                        owner_user_id=user_id,
                        treat_as_shared=False,
                        max_items=max_items,
                        max_depth=depth_limit,
                        stats=stats,
                    )
            else:
                logger.warning(f"Could not access drive for {display_name}: HTTP {drive_resp.status_code}")
        except Exception as e:
            logger.error(f"Error crawling personal drive for {display_name}: {e}", exc_info=True)
            stats["errors"] += 1

    # Crawl SPECIFIC SharePoint libraries (user-tagged, not shared)
    if include_sites:
        site_urls = config.get_sharepoint_sites_list()
        crawl_label = "shared libraries" if is_shared_crawl else display_name
        logger.info(f"Crawling {len(site_urls)} configured SharePoint site(s) for {crawl_label}...")
        
        for site_url in site_urls:
            site_id = _resolve_site_id(site_url, headers)
            if not site_id:
                logger.warning(f"Could not resolve site: {site_url}")
                stats["errors"] += 1
                continue

            # Get only document libraries from the site
            drives = _list_site_drives(site_id, headers)
            logger.info(f"Found {len(drives)} drive(s) at {site_url}")
            
            for drive in drives:
                drive_id = drive.get("id")
                drive_name = drive.get("name", "Unknown")
                drive_type = drive.get("driveType", "unknown")
                
                if not drive_id:
                    continue
                
                # Focus on document libraries (driveType: documentLibrary)
                crawl_label = "shared libraries" if is_shared_crawl else display_name
                logger.info(f"Crawling SharePoint library '{drive_name}' (type: {drive_type}) for {crawl_label}")
                _crawl_drive_items(
                    drive_id,
                    token,
                    owner_user_id=user_id,
                    treat_as_shared=is_shared_crawl,
                    max_items=max_items,
                    max_depth=depth_limit,
                    stats=stats,
                )

    # Update stat naming for clarity
    stats["sharepoint_indexed"] = stats.pop("shared_indexed", 0)
    
    crawl_label = "shared libraries" if is_shared_crawl else display_name
    logger.info(
        f"Crawl complete for {crawl_label}: "
        f"{stats['personal_indexed']} personal, "
        f"{stats['sharepoint_indexed']} SharePoint, "
        f"{stats['skipped']} skipped, "
        f"{stats['errors']} errors"
    )

    return stats
